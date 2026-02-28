import streamlit as st
import pandas as pd
import json
import re
import traceback
from anthropic import Anthropic
import os
from dotenv import load_dotenv
from datetime import datetime
import plotly.express as px
import plotly.graph_objects as go
import numpy as np
import io
import base64
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image as RLImage, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib import colors as rl_colors
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import time
try:
    import kaleido
    KALEIDO_AVAILABLE = True
except ImportError:
    KALEIDO_AVAILABLE = False

load_dotenv()

st.set_page_config(page_title="DashAI — AI Dashboard Builder", page_icon="📊", layout="wide")

st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700;800;900&display=swap');
html, body, [class*="css"] { font-family: 'Inter', -apple-system, BlinkMacSystemFont, sans-serif; }
#MainMenu {visibility: hidden;} footer {visibility: hidden;} header {visibility: hidden;}
.stApp { background: #0a0b14; min-height: 100vh; color: #e2e8f0 !important; }
.stApp p, .stApp span, .stApp div, .stApp label, .stApp li, .stApp a, .stApp small { color: #cbd5e0; }
.main .block-container { background: transparent; padding: 1.5rem 2rem 3rem 2rem; max-width: 1300px; }
h1 { font-size: 2.4rem !important; font-weight: 800 !important; letter-spacing: -0.5px; color: #f1f5f9 !important; margin-bottom: 0.2rem !important; line-height: 1.2 !important; }
h2 { color: #e2e8f0 !important; font-weight: 700 !important; font-size: 1.35rem !important; margin-top: 0 !important; margin-bottom: 0.25rem !important; letter-spacing: -0.2px; }
h3 { color: #cbd5e0 !important; font-weight: 600 !important; font-size: 1.05rem !important; margin-bottom: 0.5rem !important; }
[data-testid="metric-container"] { background: #12131f !important; border: 1px solid #1e2035 !important; border-radius: 14px !important; padding: 1.1rem 1.3rem !important; transition: border-color 0.2s ease, box-shadow 0.2s ease !important; }
[data-testid="metric-container"]:hover { border-color: #6366f1 !important; box-shadow: 0 0 0 1px #6366f140 !important; }
[data-testid="metric-container"] [data-testid="stMetricLabel"] { color: #64748b !important; font-size: 0.75rem !important; font-weight: 600 !important; text-transform: uppercase; letter-spacing: 0.08em; }
[data-testid="metric-container"] [data-testid="stMetricValue"] { color: #f8fafc !important; font-size: 1.7rem !important; font-weight: 700 !important; letter-spacing: -0.5px; }
[data-testid="metric-container"] [data-testid="stMetricDelta"] { font-size: 0.8rem !important; }
.stButton > button { background: #1e2035 !important; color: #e2e8f0 !important; border: 1px solid #2d3050 !important; border-radius: 10px !important; font-weight: 500 !important; font-size: 0.875rem !important; padding: 0.5rem 1.1rem !important; transition: all 0.15s ease !important; letter-spacing: 0.01em; }
.stButton > button:hover { background: #262840 !important; border-color: #6366f1 !important; color: #fff !important; box-shadow: 0 0 0 1px #6366f130 !important; }
.stButton > button[kind="primary"] { background: #6366f1 !important; border-color: #6366f1 !important; color: #fff !important; font-weight: 600 !important; box-shadow: 0 2px 12px #6366f150 !important; }
.stButton > button[kind="primary"]:hover { background: #4f51d1 !important; border-color: #4f51d1 !important; box-shadow: 0 4px 20px #6366f170 !important; }
[data-testid="stFileUploader"] { background: #12131f !important; border: 1.5px dashed #2d3050 !important; border-radius: 14px !important; padding: 1.5rem !important; transition: border-color 0.2s ease !important; }
[data-testid="stFileUploader"]:hover { border-color: #6366f1 !important; }
.stTextInput > div > div > input, .stTextArea > div > div > textarea { background: #12131f !important; border: 1px solid #2d3050 !important; border-radius: 10px !important; color: #e2e8f0 !important; font-size: 0.9rem !important; padding: 0.65rem 1rem !important; transition: border-color 0.15s ease !important; }
.stTextInput > div > div > input:focus, .stTextArea > div > div > textarea:focus { border-color: #6366f1 !important; box-shadow: 0 0 0 2px #6366f125 !important; }
.stSelectbox > div > div, .stMultiSelect > div > div { background: #12131f !important; border: 1px solid #2d3050 !important; border-radius: 10px !important; color: #e2e8f0 !important; }
.stMultiSelect [data-baseweb="tag"] { background: #6366f130 !important; border: 1px solid #6366f160 !important; border-radius: 6px !important; color: #a5b4fc !important; }
.stInfo { background: #1e2035 !important; border: 1px solid #3b3f6e !important; border-radius: 10px !important; color: #a5b4fc !important; }
.stSuccess { background: #0d2318 !important; border: 1px solid #166534 !important; border-radius: 10px !important; color: #86efac !important; }
.stWarning { background: #1f1707 !important; border: 1px solid #713f12 !important; border-radius: 10px !important; color: #fde68a !important; }
.stError { background: #1f0f0f !important; border: 1px solid #7f1d1d !important; border-radius: 10px !important; color: #fca5a5 !important; }
.streamlit-expanderHeader { background: #12131f !important; border: 1px solid #1e2035 !important; border-radius: 10px !important; font-weight: 600 !important; font-size: 0.875rem !important; color: #94a3b8 !important; }
.streamlit-expanderHeader:hover { border-color: #6366f1 !important; color: #e2e8f0 !important; }
[data-testid="stSidebar"] { background: #0d0e1a !important; border-right: 1px solid #1e2035 !important; }
[data-testid="stSidebar"] * { color: #cbd5e0 !important; }
[data-testid="stSidebar"] h1, [data-testid="stSidebar"] h2, [data-testid="stSidebar"] h3 { color: #f1f5f9 !important; }
.stDataFrame { border-radius: 10px !important; overflow: hidden !important; border: 1px solid #1e2035 !important; }
.stProgress > div > div > div > div { background: linear-gradient(90deg, #6366f1, #8b5cf6) !important; border-radius: 10px !important; }
hr { border: none !important; border-top: 1px solid #1e2035 !important; margin: 1.25rem 0 !important; }
.stCaption { color: #94a3b8 !important; font-size: 0.8rem !important; }
.js-plotly-plot { border-radius: 12px !important; overflow: hidden !important; }
.step-badge { display: inline-flex; align-items: center; gap: 0.4rem; background: #6366f115; border: 1px solid #6366f130; border-radius: 6px; padding: 0.25rem 0.75rem; font-size: 0.7rem; font-weight: 700; color: #818cf8; text-transform: uppercase; letter-spacing: 0.1em; margin-bottom: 0.4rem; }
.analyst-card { background: #12131f; border: 1px solid #1e2035; border-radius: 14px; padding: 1.25rem 1.5rem; margin: 0.6rem 0; }
::-webkit-scrollbar { width: 5px; height: 5px; }
::-webkit-scrollbar-track { background: #0d0e1a; }
::-webkit-scrollbar-thumb { background: #2d3050; border-radius: 10px; }
::-webkit-scrollbar-thumb:hover { background: #6366f1; }
.stTabs [data-baseweb="tab-list"] { background: #12131f !important; border-radius: 10px !important; border: 1px solid #1e2035 !important; padding: 4px !important; gap: 2px !important; }
.stTabs [data-baseweb="tab"] { border-radius: 7px !important; color: #64748b !important; font-weight: 500 !important; font-size: 0.875rem !important; padding: 0.4rem 1rem !important; }
.stTabs [aria-selected="true"] { background: #6366f1 !important; color: #fff !important; font-weight: 600 !important; }
.stCheckbox label { color: #cbd5e0 !important; font-size: 0.875rem !important; }
.stRadio label { color: #cbd5e0 !important; font-size: 0.875rem !important; }
.stMarkdown, .stMarkdown p, .stMarkdown span, .stMarkdown li, .stMarkdown a { color: #cbd5e0 !important; }
[data-testid="stFileUploader"] * { color: #94a3b8 !important; }
[data-testid="stText"], .stText { color: #cbd5e0 !important; }
</style>
""", unsafe_allow_html=True)

# ─────────────────────────────────────────────────────────
# SESSION STATE
# ─────────────────────────────────────────────────────────
for key, default in [
    ('dashboard_result', None),
    ('uploaded_df', None),
    ('cleaned_df', None),
    ('query_history', []),
    ('filters', {}),
    ('nl_results', []),
    ('saved_dashboards', []),
    ('prompt_template', ''),
    ('data_quality_issues', []),
    ('analyst_report', None),
    ('chart_style', {'color_scheme': 'Vivid', 'show_gridlines': False, 'chart_height': 420}),
    ('business_question', None),
    ('bq_confirmed', False),
    ('quality_before', None),
    ('quality_after', None),
    ('analysis_engine_result', None),
    ('recommendations_result', None),
    ('quick_wins_result', None),
    ('monitor_snapshots', []),
    ('monitor_alerts', []),
    ('monitor_kpi_config', []),
    ('monitor_ai_result', None),
    ('iteration_log', []),
]:
    if key not in st.session_state:
        st.session_state[key] = default


def compute_quality_score(df):
    score = 100
    missing_pct = df.isnull().mean().mean() * 100
    score -= min(30, missing_pct * 2)
    dup_pct = df.duplicated().mean() * 100
    score -= min(20, dup_pct * 5)
    obj_cols = df.select_dtypes(include='object').columns
    type_issues = 0
    for col in obj_cols:
        sample = df[col].dropna().head(50)
        try:
            pd.to_numeric(sample)
            type_issues += 1
        except Exception:
            pass
    score -= min(20, type_issues * 4)
    num_cols = df.select_dtypes(include='number').columns
    outlier_cols = 0
    for col in num_cols[:8]:
        q1, q3 = df[col].quantile(0.25), df[col].quantile(0.75)
        iqr = q3 - q1
        if iqr > 0:
            outlier_pct = ((df[col] < q1 - 3*iqr) | (df[col] > q3 + 3*iqr)).mean() * 100
            if outlier_pct > 5:
                outlier_cols += 1
    score -= min(15, outlier_cols * 3)
    fill_rate = (1 - df.isnull().mean().mean()) * 100
    if fill_rate > 95:
        score += 5
    return max(0, min(100, round(score)))


@st.cache_resource
def get_claude_client():
    # Try Streamlit secrets first (for Streamlit Cloud deployment)
    api_key = st.secrets.get("ANTHROPIC_API_KEY", None) if hasattr(st, 'secrets') else None
    # Fall back to environment variable (for local development)
    if not api_key:
        api_key = os.getenv('ANTHROPIC_API_KEY')
    if not api_key:
        st.error("❌ ANTHROPIC_API_KEY not found. For local dev: add to .env. For Streamlit Cloud: add to app secrets.")
        return None
    return Anthropic(api_key=api_key)

# ─────────────────────────────────────────────────────────
# COLORS
# ─────────────────────────────────────────────────────────
VIVID_COLORS = [
    "#6366f1", "#8b5cf6", "#06b6d4", "#10b981", "#f59e0b",
    "#ef4444", "#ec4899", "#3b82f6", "#84cc16", "#f97316"
]


def style_fig(fig, height=420):
    fig.update_layout(
        height=height,
        margin=dict(l=16, r=16, t=44, b=16),
        paper_bgcolor="rgba(0,0,0,0)",
        plot_bgcolor="#0d0e1a",
        font=dict(family="Inter, -apple-system, sans-serif", color="#94a3b8", size=12),
        title_font=dict(size=14, color="#e2e8f0", family="Inter"),
        legend=dict(bgcolor="#12131f", bordercolor="#1e2035", borderwidth=1, font=dict(color="#94a3b8", size=11)),
        colorway=VIVID_COLORS,
        xaxis=dict(gridcolor="#1e2035", linecolor="#1e2035", tickcolor="#1e2035", tickfont=dict(color="#64748b", size=11), title_font=dict(color="#94a3b8"), showgrid=False, zeroline=False),
        yaxis=dict(gridcolor="#1e2035", linecolor="rgba(0,0,0,0)", tickcolor="#1e2035", tickfont=dict(color="#64748b", size=11), title_font=dict(color="#94a3b8"), showgrid=True, zeroline=False),
    )
    fig.update_traces(marker_line_width=0, selector=dict(type="bar"))
    traces = fig.data
    if len(traces) == 1:
        tr = traces[0]
        ttype = tr.type if hasattr(tr, 'type') else ''
        if ttype == 'bar':
            x_vals = tr.x if tr.x is not None else []
            n = len(x_vals)
            if n > 1:
                fig.update_traces(marker=dict(color=[VIVID_COLORS[i % len(VIVID_COLORS)] for i in range(n)]))
        elif ttype == 'histogram':
            fig.update_traces(marker=dict(color=VIVID_COLORS[2]))
    elif len(traces) > 1:
        for i, tr in enumerate(traces):
            ttype = tr.type if hasattr(tr, 'type') else ''
            if ttype in ('bar', 'scatter', 'box', 'violin'):
                color = VIVID_COLORS[i % len(VIVID_COLORS)]
                try:
                    fig.data[i].marker.color = color
                    if ttype == 'scatter' and tr.mode and 'lines' in tr.mode:
                        fig.data[i].line.color = color
                except Exception:
                    pass
    return fig


# ─────────────────────────────────────────────────────────
# STATE CODES
# ─────────────────────────────────────────────────────────
STATE_CODES = {
    'Alabama':'AL','Alaska':'AK','Arizona':'AZ','Arkansas':'AR','California':'CA',
    'Colorado':'CO','Connecticut':'CT','Delaware':'DE','Florida':'FL','Georgia':'GA',
    'Hawaii':'HI','Idaho':'ID','Illinois':'IL','Indiana':'IN','Iowa':'IA',
    'Kansas':'KS','Kentucky':'KY','Louisiana':'LA','Maine':'ME','Maryland':'MD',
    'Massachusetts':'MA','Michigan':'MI','Minnesota':'MN','Mississippi':'MS','Missouri':'MO',
    'Montana':'MT','Nebraska':'NE','Nevada':'NV','New Hampshire':'NH','New Jersey':'NJ',
    'New Mexico':'NM','New York':'NY','North Carolina':'NC','North Dakota':'ND','Ohio':'OH',
    'Oklahoma':'OK','Oregon':'OR','Pennsylvania':'PA','Rhode Island':'RI','South Carolina':'SC',
    'South Dakota':'SD','Tennessee':'TN','Texas':'TX','Utah':'UT','Vermont':'VT',
    'Virginia':'VA','Washington':'WA','West Virginia':'WV','Wisconsin':'WI','Wyoming':'WY'
}
STATE_NAMES = {v: k for k, v in STATE_CODES.items()}
ALL_STATE_CODES = list(STATE_CODES.values())


def make_us_map(df_agg, value_col, state_col='state', title='', colorscale='Blues', show_all_states=True):
    def to_code(s):
        s = str(s).strip()
        if len(s) == 2 and s.upper() in ALL_STATE_CODES:
            return s.upper()
        return STATE_CODES.get(s.title(), STATE_CODES.get(s, None))
    df_map = df_agg.copy()
    df_map['_code'] = df_map[state_col].apply(to_code)
    df_map = df_map[df_map['_code'].notna()].copy()
    if show_all_states:
        existing_codes = set(df_map['_code'].tolist())
        missing_codes = [c for c in ALL_STATE_CODES if c not in existing_codes]
        if missing_codes:
            padding = pd.DataFrame({'_code': missing_codes, value_col: [None] * len(missing_codes)})
            df_map = pd.concat([df_map, padding], ignore_index=True)
    fig = go.Figure(data=go.Choropleth(
        locations=df_map['_code'], z=df_map[value_col], locationmode='USA-states',
        colorscale=colorscale, autocolorscale=False, marker_line_color='#1e2035', marker_line_width=0.5,
        colorbar=dict(title=dict(text=value_col.replace('_', ' ').title(), font=dict(color='#94a3b8', size=11)),
                      tickfont=dict(color='#64748b', size=10), bgcolor='#12131f', bordercolor='#1e2035', len=0.8, thickness=14),
    ))
    fig.update_layout(
        title_text=title,
        geo=dict(scope='usa', projection_type='albers usa', showland=True, landcolor='#12131f',
                 showlakes=False, lakecolor='#0d0e1a', showframe=False, bgcolor='rgba(0,0,0,0)',
                 coastlinecolor='#2d3050', showcoastlines=True),
    )
    return fig


# ─────────────────────────────────────────────────────────
# UTILITIES
# ─────────────────────────────────────────────────────────
def extract_json(text):
    text = text.strip()
    text = re.sub(r'^```(?:json)?\s*\n?', '', text)
    text = re.sub(r'\n?```\s*$', '', text).strip()
    try:
        return json.loads(text)
    except json.JSONDecodeError:
        pass
    brace_start = text.find('{')
    if brace_start == -1:
        return None
    depth, brace_end = 0, -1
    in_string, escape_next = False, False
    for i in range(brace_start, len(text)):
        ch = text[i]
        if escape_next: escape_next = False; continue
        if ch == '\\': escape_next = True; continue
        if ch == '"': in_string = not in_string; continue
        if in_string: continue
        if ch == '{': depth += 1
        elif ch == '}':
            depth -= 1
            if depth == 0: brace_end = i; break
    if brace_end == -1:
        return None
    json_str = text[brace_start:brace_end + 1]
    try:
        return json.loads(json_str)
    except json.JSONDecodeError:
        pass
    cleaned = re.sub(r',\s*([}\]])', r'\1', json_str)
    try:
        return json.loads(cleaned)
    except json.JSONDecodeError:
        return None


def safe_exec(code, df):
    code = re.sub(r'if\s+(\w+)\s*(?:==\s*None|is None)\s*:', r'if \1 is None:', code)
    code = re.sub(r'\bif\s+((?:df|df_)\w*)\s*:', r'if \1 is not None and not \1.empty:', code)
    code = re.sub(r'\bif\s+not\s+((?:df|df_)\w*)\s*:', r'if \1 is None or \1.empty:', code)
    code = re.sub(r'\band\s+((?:df|df_)\w*)\s*:', r'and \1 is not None and not \1.empty:', code)
    exec_globals = {
        'df': df.copy(), 'pd': pd, 'px': px, 'go': go, 'np': np,
        'state_codes': STATE_CODES, 'STATE_CODES': STATE_CODES,
        'ALL_STATE_CODES': ALL_STATE_CODES, 'make_us_map': make_us_map,
        'VIVID_COLORS': VIVID_COLORS, 'result': None, 'fig': None, 'insight': '',
    }
    try:
        exec(code, exec_globals)
        return {'result': exec_globals.get('result') or exec_globals.get('fig'), 'insight': exec_globals.get('insight', ''), 'error': None}
    except Exception as e:
        error_msg = str(e)
        if 'unexpected keyword argument' in error_msg:
            bad_kwarg = re.search(r"unexpected keyword argument '(\w+)'", error_msg)
            if bad_kwarg:
                fixed_code = re.sub(rf',?\s*{bad_kwarg.group(1)}=[^,)]+', '', code)
                try:
                    eg2 = {'df': df.copy(), 'pd': pd, 'px': px, 'go': go, 'np': np,
                           'state_codes': STATE_CODES, 'VIVID_COLORS': VIVID_COLORS,
                           'make_us_map': make_us_map, 'ALL_STATE_CODES': ALL_STATE_CODES,
                           'result': None, 'fig': None, 'insight': ''}
                    exec(fixed_code, eg2)
                    return {'result': eg2.get('result') or eg2.get('fig'), 'insight': eg2.get('insight', ''), 'error': None}
                except Exception:
                    pass
        return {'result': None, 'insight': '', 'error': error_msg}


def validate_data(df):
    issues = []
    missing_pct = (df.isnull().sum() / len(df) * 100).round(2)
    high_missing = missing_pct[missing_pct > 20]
    if len(high_missing) > 0:
        issues.append({'severity': 'warning', 'type': 'Missing Values', 'message': f"High missing values (>20%) in: {', '.join(high_missing.index)}"})
    dupes = df.duplicated().sum()
    if dupes > 0:
        issues.append({'severity': 'warning', 'type': 'Duplicates', 'message': f"{dupes:,} duplicate rows found ({dupes/len(df)*100:.1f}%)"})
    numeric_cols = df.select_dtypes(include=['number']).columns
    for col in numeric_cols[:5]:
        Q1, Q3 = df[col].quantile(0.25), df[col].quantile(0.75)
        IQR = Q3 - Q1
        outliers = ((df[col] < (Q1 - 1.5 * IQR)) | (df[col] > (Q3 + 1.5 * IQR))).sum()
        if outliers > len(df) * 0.05:
            issues.append({'severity': 'info', 'type': 'Outliers', 'message': f"{outliers} potential outliers in '{col}' ({outliers/len(df)*100:.1f}%)"})
    return issues


# ─────────────────────────────────────────────────────────
# AI ANALYST
# ─────────────────────────────────────────────────────────
def run_ai_analyst(df, client, business_question=None):
    profile = {
        "shape": list(df.shape), "columns": list(df.columns),
        "dtypes": {col: str(dtype) for col, dtype in df.dtypes.items()},
        "missing_pct": {col: round(df[col].isnull().sum() / len(df) * 100, 2) for col in df.columns},
        "unique_counts": {col: int(df[col].nunique()) for col in df.columns},
        "numeric_stats": {},
        "sample_values": {col: df[col].dropna().unique()[:8].tolist() for col in df.columns[:20]},
        "duplicate_rows": int(df.duplicated().sum()),
    }
    for col in df.select_dtypes(include=['number']).columns[:15]:
        profile["numeric_stats"][col] = {
            "min": float(df[col].min()), "max": float(df[col].max()),
            "mean": float(df[col].mean()), "median": float(df[col].median()),
            "std": float(df[col].std()), "nulls": int(df[col].isnull().sum())
        }
    prompt = f"""You are a Senior Data Analyst. Analyze this dataset and produce a complete data cleaning + transformation plan.

DATASET PROFILE:
{json.dumps(profile, indent=2, default=str)}

Respond with ONLY this JSON (no other text):
{{
  "data_summary": {{"title": "Dataset name/type", "overview": "2-3 sentence overview", "quality_score": 85, "key_findings": ["finding1", "finding2", "finding3"]}},
  "issues_found": [{{"issue_type": "Missing Values / Duplicates / Wrong Type / Outliers / Inconsistent Format / Bad Column Names / Percentage Error", "column": "column_name or ALL", "severity": "critical / high / medium / low", "description": "Specific description", "fix_description": "What will be done"}}],
  "cleaning_code": "# Complete Python pandas cleaning code\\nimport pandas as pd\\nimport numpy as np\\n\\ndf_clean = df.copy()\\n# ... all transformations\\nresult = df_clean",
  "engineered_features": [{{"feature_name": "new_column_name", "description": "What this feature represents", "code": "df_clean['new_col'] = ..."}}],
  "analysis_insights": [{{"category": "Revenue / Customer / Operations / Risk / Growth", "insight": "Specific data-driven insight with numbers", "recommendation": "Actionable recommendation"}}],
  "dashboard_suggestions": ["Suggestion 1", "Suggestion 2", "Suggestion 3"]
}}

BUSINESS CONTEXT: {f"Focus all insights and recommendations on answering: '{business_question}'" if business_question else "No specific business question set — provide general insights."}

CLEANING RULES:
1. Fix ALL missing values (impute with median/mode/forward-fill, never just drop unless >50% missing)
2. Remove duplicate rows
3. Fix data types (convert string numbers to float, parse dates)
4. Normalize percentage columns that are clearly wrong (>100% for rates, >1000% is a bug)
5. Fix negative values where they shouldn't exist
6. Standardize categorical values (strip whitespace, consistent casing)
7. Engineer 2-3 useful derived features (e.g. ARPU, churn_risk_score, revenue_tier)
8. cleaning_code must be complete executable Python, uses 'df' as input, assigns final cleaned dataframe to 'result'
9. CRITICAL CODE RULES:
   - NEVER use pd.cut() or pd.qcut() with labels parameter — use np.select() instead
   - NEVER write bare 'if dataframe_variable:' — always: if len(df_clean) > 0:
   - NEVER use chained assignment — always use .loc[]
   - ALWAYS convert date columns with pd.to_datetime() BEFORE using .dt accessor"""
    try:
        response = client.messages.create(model="claude-sonnet-4-20250514", max_tokens=6000, messages=[{"role": "user", "content": prompt}])
        return extract_json(response.content[0].text)
    except Exception as e:
        st.error(f"AI Analyst error: {e}")
        return None


# ─────────────────────────────────────────────────────────
# DATETIME AUTO-FIX HELPERS
# ─────────────────────────────────────────────────────────
def _preprocess_datetime_cols(df):
    """Auto-convert columns that look like dates/times but stored as object strings."""
    df = df.copy()
    date_pat = re.compile(
        r'^\d{4}[-/]\d{2}[-/]\d{2}'
        r'|^\d{2}[-/]\d{2}[-/]\d{4}'
        r'|^\d{4}\d{2}\d{2}$'
        r'|^\d{4}-\d{2}-\d{2}T'
    )
    for col in df.select_dtypes(include=['object']).columns:
        sample = df[col].dropna().head(50).astype(str)
        if sample.empty:
            continue
        n_datelike = sample.apply(lambda x: bool(date_pat.match(x))).sum()
        if n_datelike >= len(sample) * 0.6:
            try:
                df[col] = pd.to_datetime(df[col], infer_datetime_format=True, errors='coerce')
            except Exception:
                pass
    return df


def _fix_dt_accessor_in_code(cleaning_code):
    """Inject pd.to_datetime() for any column used with .dt accessor."""
    dt_cols = re.findall(r"df(?:_clean|_cleaned|_final|_qw)?\['([^']+)'\]\.dt\.", cleaning_code)
    dt_cols += re.findall(r'df(?:_clean|_cleaned|_final|_qw)?\["([^"]+)"\]\.dt\.', cleaning_code)
    if not dt_cols:
        return None
    convert_lines = [
        f"df_clean['{col}'] = pd.to_datetime(df_clean['{col}'], infer_datetime_format=True, errors='coerce')"
        for col in set(dt_cols)
    ]
    lines = cleaning_code.split('\n')
    insert_idx = 0
    for i, line in enumerate(lines):
        stripped = line.strip()
        if ('df_clean' in stripped or 'df_cleaned' in stripped) and 'copy()' in stripped:
            insert_idx = i + 1
            break
    lines = lines[:insert_idx] + convert_lines + lines[insert_idx:]
    return '\n'.join(lines)


def _safe_clean_exec(cleaning_code, df):
    # PRE-PROCESS: auto-convert date-like string columns so .dt accessor works
    df = _preprocess_datetime_cols(df)

    # Fix bare DataFrame truth checks
    cleaning_code = re.sub(r'\bif\s+(\w*df\w*|df_\w*|\w*_df)\s*:', r'if len(\1) > 0:', cleaning_code)
    cleaning_code = re.sub(r'\bif\s+not\s+(\w*df\w*|df_\w*|\w*_df)\s*:', r'if len(\1) == 0:', cleaning_code)
    cleaning_code = re.sub(r'\bor\s+(\w*df\w*|df_\w*|\w*_df)\b(?!\s*\.)', r'or len(\1) > 0', cleaning_code)
    cleaning_code = re.sub(r'\band\s+(\w*df\w*|df_\w*|\w*_df)\b(?!\s*\.)', r'and len(\1) > 0', cleaning_code)
    cleaning_code = cleaning_code.replace('.copy(False)', '.copy()')

    base_ns = {'df': df.copy(), 'pd': pd, 'np': np, 'result': None}

    def _try_exec(code):
        ns = {**base_ns}
        exec(code, ns)
        result = ns.get('result')
        if result is None:
            for key in ['df_clean', 'df_cleaned', 'cleaned_df', 'df_final']:
                if key in ns and isinstance(ns[key], pd.DataFrame):
                    result = ns[key]
                    break
        return result

    try:
        result = _try_exec(cleaning_code)
        if result is None:
            return {'result': None, 'error': 'No result dataframe found.'}
        return {'result': result, 'error': None}

    except Exception as e:
        err = str(e)

        # AUTO-FIX 1: .dt accessor on non-datetime column
        if 'dt accessor' in err.lower() or 'datetimelike' in err.lower():
            fixed_code = _fix_dt_accessor_in_code(cleaning_code)
            if fixed_code:
                try:
                    result = _try_exec(fixed_code)
                    if result is not None:
                        return {'result': result, 'error': None}
                except Exception as e2:
                    err = str(e2)

        # AUTO-FIX 2: truth value of DataFrame ambiguous
        if 'truth value' in err.lower() and 'ambiguous' in err.lower():
            try:
                fixed = re.sub(
                    r'\bif\s+(\w+)\s*:',
                    lambda m: f'if len({m.group(1)}) > 0:'
                    if m.group(1) not in ('True', 'False', 'None', 'not', 'len', 'isinstance', 'type')
                    else m.group(0),
                    cleaning_code
                )
                result = _try_exec(fixed)
                if result is not None:
                    return {'result': result, 'error': None}
            except Exception:
                pass

        return {'result': None, 'error': err}


def render_analyst_report(report, df, client):
    summary = report.get("data_summary", {})
    issues = report.get("issues_found", [])
    insights = report.get("analysis_insights", [])
    features = report.get("engineered_features", [])
    suggestions = report.get("dashboard_suggestions", [])
    quality_score = summary.get("quality_score", 0)

    st.markdown("### Data Quality Score")
    score_color = "#10b981" if quality_score >= 80 else "#f59e0b" if quality_score >= 60 else "#ef4444"
    fig_gauge = go.Figure(go.Indicator(
        mode="gauge+number", value=quality_score,
        title={"text": "Overall Quality", "font": {"color": "#94a3b8", "size": 13}},
        number={"suffix": "%", "font": {"color": score_color, "size": 34}},
        gauge={"axis": {"range": [0, 100], "tickcolor": "#2d3050", "tickfont": {"color": "#64748b"}},
               "bar": {"color": score_color, "thickness": 0.25}, "bgcolor": "#12131f", "borderwidth": 0,
               "steps": [{"range": [0, 60], "color": "#1f0f0f"}, {"range": [60, 80], "color": "#1f1707"}, {"range": [80, 100], "color": "#0d2318"}]}
    ))
    fig_gauge.update_layout(height=200, margin=dict(l=20, r=20, t=40, b=10), paper_bgcolor="rgba(0,0,0,0)", font=dict(color="#94a3b8", family="Inter"))
    col1, col2 = st.columns([1, 2])
    with col1:
        st.plotly_chart(fig_gauge, use_container_width=True)
    with col2:
        st.markdown(f"**{summary.get('title', 'Dataset')}**")
        bq_active = st.session_state.get('business_question')
        if bq_active:
            st.markdown(f"<span style='color:#6366f1;font-size:0.78rem;font-weight:600'>🎯 Analyzing: {bq_active}</span>", unsafe_allow_html=True)
        st.markdown(f"<span style='color:#64748b;font-size:0.875rem'>{summary.get('overview', '')}</span>", unsafe_allow_html=True)
        st.markdown("**Key Findings:**")
        for f in summary.get("key_findings", []):
            st.markdown(f"<span style='color:#6366f1;font-size:0.875rem'>↳</span> <span style='color:#94a3b8;font-size:0.875rem'>{f}</span>", unsafe_allow_html=True)

    st.divider()

    if issues:
        st.markdown("### Issues Detected")
        severity_border = {"critical": "#ef4444", "high": "#f97316", "medium": "#f59e0b", "low": "#10b981"}
        critical = [i for i in issues if i.get('severity') == 'critical']
        high = [i for i in issues if i.get('severity') == 'high']
        medium = [i for i in issues if i.get('severity') == 'medium']
        low = [i for i in issues if i.get('severity') == 'low']
        ic1, ic2, ic3, ic4 = st.columns(4)
        ic1.metric("Critical", len(critical)); ic2.metric("High", len(high)); ic3.metric("Medium", len(medium)); ic4.metric("Low", len(low))
        with st.expander("View All Issues & Fixes", expanded=True):
            for issue in issues:
                sev = issue.get('severity', 'low')
                border_color = severity_border.get(sev, '#64748b')
                st.markdown(f"""<div style='background:#12131f;border-radius:8px;padding:0.75rem 1rem;margin:0.35rem 0;border-left:2px solid {border_color}'>
                <span style='color:{border_color};font-size:0.7rem;font-weight:700;text-transform:uppercase;letter-spacing:0.08em'>{sev}</span>
                &nbsp;·&nbsp;<strong style='color:#e2e8f0;font-size:0.875rem'>{issue.get("issue_type","Issue")}</strong>
                &nbsp;<code style='color:#6366f1;font-size:0.8rem;background:#6366f115;padding:0.1rem 0.4rem;border-radius:4px'>{issue.get("column","")}</code><br>
                <span style='color:#64748b;font-size:0.825rem'>{issue.get("description","")}</span><br>
                <span style='color:#10b981;font-size:0.8rem'>✓ {issue.get("fix_description","")}</span></div>""", unsafe_allow_html=True)

    st.divider()

    if features:
        st.markdown("### Features to Engineer")
        for feat in features:
            st.markdown(f"""<div style='background:#12131f;border-radius:8px;padding:0.7rem 1rem;margin:0.3rem 0;border:1px solid #1e2035'>
            <strong style='color:#8b5cf6;font-size:0.875rem'>+ {feat.get("feature_name","")}</strong><br>
            <span style='color:#64748b;font-size:0.825rem'>{feat.get("description","")}</span></div>""", unsafe_allow_html=True)

    st.divider()

    if insights:
        st.markdown("### Business Insights")
        cat_colors = {"Revenue": "#6366f1", "Customer": "#8b5cf6", "Operations": "#06b6d4", "Risk": "#ef4444", "Growth": "#10b981"}
        for ins in insights:
            cat = ins.get("category", "General")
            color = cat_colors.get(cat, "#64748b")
            st.markdown(f"""<div style='background:#12131f;border-radius:8px;padding:0.9rem 1rem;margin:0.4rem 0;border-left:2px solid {color}'>
            <span style='color:{color};font-size:0.7rem;font-weight:700;text-transform:uppercase;letter-spacing:0.08em'>{cat}</span><br>
            <strong style='color:#e2e8f0;font-size:0.875rem'>{ins.get("insight","")}</strong><br>
            <span style='color:#10b981;font-size:0.825rem'>→ {ins.get("recommendation","")}</span></div>""", unsafe_allow_html=True)

    st.divider()

    st.markdown("### Apply AI Cleaning & Transformation")
    cleaning_code = report.get("cleaning_code", "")
    if cleaning_code:
        with st.expander("View Generated Cleaning Code"):
            st.code(cleaning_code, language="python")
        col1, col2 = st.columns(2)
        with col1:
            if st.button("✨ Apply All Fixes & Engineer Features", type="primary", use_container_width=True, key="btn_apply_all_fixes"):
                with st.spinner("Cleaning and transforming data..."):
                    raw_df = st.session_state.get('uploaded_df', df)
                    out = _safe_clean_exec(cleaning_code, raw_df)
                    if out['error']:
                        st.error(f"Cleaning error: {out['error']}")
                    elif out['result'] is not None and isinstance(out['result'], pd.DataFrame):
                        cleaned = out['result']
                        st.session_state['cleaned_df'] = cleaned
                        st.session_state['dashboard_result'] = None
                        st.session_state['quality_after'] = compute_quality_score(cleaned)
                        st.rerun()
                    else:
                        st.warning("Cleaning returned no result. Check code above.")
        with col2:
            if st.button("📥 Download Cleaned Data", use_container_width=True, key="btn_download_cleaned"):
                raw_df = st.session_state.get('uploaded_df', df)
                out = _safe_clean_exec(cleaning_code, raw_df)
                if out['result'] is not None and isinstance(out['result'], pd.DataFrame):
                    csv = out['result'].to_csv(index=False)
                    st.download_button("⬇️ Download CSV", csv, "cleaned_data.csv", "text/csv", use_container_width=True, key="btn_dl_csv")

    if st.session_state.get('quality_before') is not None and st.session_state.get('quality_after') is not None:
        before = st.session_state['quality_before']
        after = st.session_state['quality_after']
        delta = after - before
        delta_color = "#10b981" if delta >= 0 else "#ef4444"
        delta_icon = "↑" if delta >= 0 else "↓"
        st.markdown(f"""<div style='background:#12131f;border:1px solid #1e2035;border-radius:14px;padding:1.25rem 1.5rem;margin:0.75rem 0'>
          <div style='font-size:0.7rem;font-weight:700;color:#6366f1;text-transform:uppercase;letter-spacing:0.1em;margin-bottom:0.75rem'>✅ Data Quality — Before vs After Cleaning</div>
          <div style='display:flex;align-items:center;gap:2rem;flex-wrap:wrap'>
            <div style='text-align:center;min-width:90px'>
              <div style='font-size:0.7rem;color:#64748b;font-weight:600;text-transform:uppercase;letter-spacing:0.06em;margin-bottom:0.25rem'>Before</div>
              <div style='font-size:2.2rem;font-weight:800;color:{"#10b981" if before>=80 else "#f59e0b" if before>=60 else "#ef4444"}'>{before}%</div>
            </div>
            <div style='text-align:center;flex:1;min-width:80px'>
              <div style='font-size:2rem;color:#2d3050'>→</div>
              <div style='font-size:1rem;font-weight:700;color:{delta_color}'>{delta_icon} {abs(delta)} pts</div>
            </div>
            <div style='text-align:center;min-width:90px'>
              <div style='font-size:0.7rem;color:#64748b;font-weight:600;text-transform:uppercase;letter-spacing:0.06em;margin-bottom:0.25rem'>After</div>
              <div style='font-size:2.2rem;font-weight:800;color:{"#10b981" if after>=80 else "#f59e0b" if after>=60 else "#ef4444"}'>{after}%</div>
            </div>
          </div>
        </div>""", unsafe_allow_html=True)

    if suggestions:
        st.divider()
        st.markdown("### Recommended Next Steps")
        for i, sug in enumerate(suggestions, 1):
            st.markdown(f"<span style='color:#6366f1;font-size:0.875rem'>{i}.</span> <span style='color:#94a3b8;font-size:0.875rem'>{sug}</span>", unsafe_allow_html=True)


# ─────────────────────────────────────────────────────────
# EXPORTS
# ─────────────────────────────────────────────────────────
def export_to_pdf(dashboard, df):
    try:
        buffer = io.BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=letter, topMargin=0.75*inch, bottomMargin=0.75*inch, leftMargin=0.75*inch, rightMargin=0.75*inch)
        story = []
        styles = getSampleStyleSheet()
        title_style = ParagraphStyle('DashTitle', parent=styles['Heading1'], fontSize=22, textColor=rl_colors.HexColor('#6366f1'), spaceAfter=8, fontName='Helvetica-Bold')
        sub_style = ParagraphStyle('DashSub', parent=styles['Normal'], fontSize=11, textColor=rl_colors.HexColor('#64748b'), spaceAfter=20)
        section_style = ParagraphStyle('Section', parent=styles['Heading2'], fontSize=14, textColor=rl_colors.HexColor('#1e293b'), spaceBefore=20, spaceAfter=10, fontName='Helvetica-Bold')
        story.append(Paragraph(dashboard.get('title', 'Dashboard'), title_style))
        story.append(Paragraph(dashboard.get('description', ''), sub_style))
        kpis = dashboard.get('kpis', [])
        if kpis:
            story.append(Paragraph("Key Metrics", section_style))
            kpi_data = [['Metric', 'Value']]
            for kpi in kpis:
                out = safe_exec(kpi.get('code', 'result=0'), df)
                val = out['result'] if not out['error'] else 'N/A'
                try: display_val = kpi.get('format', '{:,}').format(val)
                except: display_val = str(val)
                kpi_data.append([kpi.get('label', 'KPI'), display_val])
            t = Table(kpi_data, colWidths=[3.5*inch, 3*inch])
            t.setStyle(TableStyle([('BACKGROUND',(0,0),(-1,0),rl_colors.HexColor('#6366f1')),('TEXTCOLOR',(0,0),(-1,0),rl_colors.white),('FONTNAME',(0,0),(-1,0),'Helvetica-Bold'),('FONTSIZE',(0,0),(-1,0),12),('ALIGN',(0,0),(-1,-1),'CENTER'),('ROWBACKGROUNDS',(0,1),(-1,-1),[rl_colors.HexColor('#f8fafc'),rl_colors.white]),('GRID',(0,0),(-1,-1),0.5,rl_colors.HexColor('#e2e8f0')),('TOPPADDING',(0,0),(-1,-1),8),('BOTTOMPADDING',(0,0),(-1,-1),8)]))
            story.append(t)
            story.append(Spacer(1, 0.3*inch))
        insights = dashboard.get('insights', [])
        if insights:
            story.append(Paragraph("Key Insights", section_style))
            for ins in insights:
                story.append(Paragraph(f"→ {ins}", styles['Normal']))
                story.append(Spacer(1, 0.1*inch))
        vizs = dashboard.get('visualizations', [])
        if vizs:
            story.append(Paragraph("Visualizations", section_style))
            for viz in vizs:
                story.append(Paragraph(viz.get('title', 'Chart'), ParagraphStyle('ChartTitle', parent=styles['Normal'], fontSize=12, fontName='Helvetica-Bold', textColor=rl_colors.HexColor('#1e293b'), spaceAfter=6)))
                out = safe_exec(viz.get('code', ''), df)
                if not out['error'] and out['result'] is not None and hasattr(out['result'], 'update_layout'):
                    fig = out['result']
                    fig.update_layout(paper_bgcolor='white', plot_bgcolor='#f8faff', font=dict(color='#333', family='Helvetica'), height=320, width=680, margin=dict(l=40, r=40, t=40, b=40))
                    try:
                        if not KALEIDO_AVAILABLE: raise RuntimeError("kaleido not installed")
                        img_bytes = fig.to_image(format="png", scale=1.5)
                        img = RLImage(io.BytesIO(img_bytes), width=6.5*inch, height=3*inch)
                        story.append(img)
                    except Exception:
                        story.append(Paragraph(f"[Chart: {viz.get('title', '')} — install kaleido for chart images]", styles['Normal']))
                if viz.get('insight'):
                    story.append(Paragraph(f"💡 {viz['insight']}", ParagraphStyle('Insight', parent=styles['Normal'], fontSize=10, textColor=rl_colors.HexColor('#64748b'), spaceAfter=15, leftIndent=10)))
                story.append(Spacer(1, 0.2*inch))
        doc.build(story)
        buffer.seek(0)
        return buffer
    except Exception as e:
        st.error(f"PDF export error: {e}")
        return None


def export_to_powerpoint(dashboard, df):
    try:
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]

        def add_textbox(slide, text, left, top, width, height, size=18, bold=False, color="FFFFFF", align=PP_ALIGN.LEFT):
            tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
            tf = tb.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.text = text; p.alignment = align
            run = p.runs[0] if p.runs else p.add_run()
            run.font.size = Pt(size); run.font.bold = bold
            run.font.color.rgb = RGBColor.from_string(color)

        def add_rect(slide, left, top, width, height, hex_color="6366f1"):
            shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
            shape.fill.solid()
            r, g, b = int(hex_color[0:2],16), int(hex_color[2:4],16), int(hex_color[4:6],16)
            shape.fill.fore_color.rgb = RGBColor(r, g, b)
            shape.line.fill.background()

        slide = prs.slides.add_slide(blank)
        add_rect(slide, 0, 0, 13.33, 7.5, "0a0b14"); add_rect(slide, 0, 6.5, 13.33, 1.0, "6366f1")
        add_textbox(slide, dashboard.get('title','Dashboard'), 0.8, 2.2, 11.5, 1.5, size=38, bold=True, color="F1F5F9")
        add_textbox(slide, dashboard.get('description',''), 0.8, 3.8, 11, 1.2, size=17, color="94A3B8")
        add_textbox(slide, f"DashAI  ·  {datetime.now().strftime('%B %d, %Y')}", 0.8, 6.6, 10, 0.5, size=11, color="6366f1")

        kpis = dashboard.get('kpis', [])
        if kpis:
            slide = prs.slides.add_slide(blank)
            add_rect(slide, 0, 0, 13.33, 7.5, "0a0b14"); add_rect(slide, 0, 0, 13.33, 1.0, "12131f")
            add_textbox(slide, "Key Metrics", 0.5, 0.18, 10, 0.65, size=22, bold=True, color="F1F5F9")
            n = len(kpis); card_w = 11.5 / n
            for i, kpi in enumerate(kpis):
                out = safe_exec(kpi.get('code', 'result=0'), df)
                val = out['result'] if not out['error'] else 'N/A'
                try: display_val = kpi.get('format', '{:,}').format(val)
                except: display_val = str(val)
                left = 0.8 + i * (card_w + 0.15)
                add_rect(slide, left, 1.4, card_w, 2.2, "12131f")
                add_textbox(slide, kpi.get('label','KPI'), left+0.1, 1.55, card_w-0.2, 0.5, size=11, color="64748B", align=PP_ALIGN.CENTER)
                add_textbox(slide, display_val, left+0.1, 2.1, card_w-0.2, 0.9, size=26, bold=True, color="F8FAFC", align=PP_ALIGN.CENTER)

        insights = dashboard.get('insights', [])
        if insights:
            slide = prs.slides.add_slide(blank)
            add_rect(slide, 0, 0, 13.33, 7.5, "0a0b14"); add_rect(slide, 0, 0, 13.33, 1.0, "12131f")
            add_textbox(slide, "Key Insights", 0.5, 0.18, 10, 0.65, size=22, bold=True, color="F1F5F9")
            for j, ins in enumerate(insights[:5]):
                add_rect(slide, 0.6, 1.2 + j*1.1, 12, 0.9, "12131f")
                add_textbox(slide, f"  →  {ins}", 0.7, 1.3 + j*1.1, 11.5, 0.7, size=13, color="94A3B8")

        for viz in dashboard.get('visualizations', []):
            out = safe_exec(viz.get('code', ''), df)
            slide = prs.slides.add_slide(blank)
            add_rect(slide, 0, 0, 13.33, 7.5, "0a0b14"); add_rect(slide, 0, 0, 13.33, 1.0, "12131f")
            add_textbox(slide, viz.get('title','Chart'), 0.5, 0.18, 11, 0.65, size=20, bold=True, color="F1F5F9")
            if viz.get('insight'):
                add_textbox(slide, f"→ {viz['insight']}", 0.5, 6.85, 12, 0.5, size=11, color="64748B")
            if not out['error'] and out['result'] is not None and hasattr(out['result'], 'update_layout'):
                fig = out['result']
                fig.update_layout(paper_bgcolor='#0a0b14', plot_bgcolor='#12131f', font=dict(color='#94a3b8', family='Arial'), height=420, width=900, margin=dict(l=50, r=50, t=50, b=50), colorway=VIVID_COLORS)
                try:
                    if not KALEIDO_AVAILABLE: raise RuntimeError("kaleido not installed")
                    img_bytes = fig.to_image(format="png", scale=2)
                    slide.shapes.add_picture(io.BytesIO(img_bytes), Inches(1.5), Inches(1.2), Inches(10.3), Inches(5.3))
                except Exception:
                    add_textbox(slide, "[Install kaleido for chart images]", 1.5, 3.5, 10, 0.8, size=13, color="64748B", align=PP_ALIGN.CENTER)

        buffer = io.BytesIO()
        prs.save(buffer); buffer.seek(0)
        return buffer
    except Exception as e:
        st.error(f"PowerPoint export error: {e}")
        return None


# ─────────────────────────────────────────────────────────
# PROMPT TEMPLATES
# ─────────────────────────────────────────────────────────
PROMPT_TEMPLATES = {
    "Executive": {"emoji": "👔", "prompt": "Executive dashboard with revenue KPIs (total revenue, active customers, average satisfaction, churn rate), monthly revenue growth trend line chart, plan type distribution donut chart, and geographic revenue map by state"},
    "Sales":     {"emoji": "📈", "prompt": "Sales performance dashboard with pipeline KPIs, revenue by plan type bar chart, top states by revenue, sales trend over time line chart, and average revenue per user"},
    "Marketing": {"emoji": "🎯", "prompt": "Marketing analytics dashboard showing customer acquisition by plan type, customer satisfaction distribution histogram, customer segmentation scatter plot, plan mix breakdown, and churn vs satisfaction analysis"},
    "Operations":{"emoji": "⚙️", "prompt": "Telecom operations dashboard with network KPIs, data usage distribution by plan (box plot), customer count by state bar chart, service type breakdown pie chart, and monthly active customers trend. For all percentage/rate metrics take the MEAN not SUM."},
    "Finance":   {"emoji": "💰", "prompt": "Financial dashboard with revenue breakdown by plan type, average revenue per user by state, revenue distribution histogram, monthly revenue trend, and top 10 revenue states bar chart"},
    "Customer":  {"emoji": "👥", "prompt": "Customer analytics dashboard with satisfaction score distribution, churn analysis by plan type, customer tenure histogram, data usage patterns by plan, and customer count by state map"},
}


# ─────────────────────────────────────────────────────────
# NL QUERY
# ─────────────────────────────────────────────────────────
def natural_language_query(df, question, client):
    df_info = {
        "columns": list(df.columns), "dtypes": {col: str(dtype) for col, dtype in df.dtypes.items()},
        "shape": list(df.shape), "sample_values": {col: df[col].dropna().unique()[:5].tolist() for col in df.columns[:15]},
        "numeric_stats": {col: {"min": float(df[col].min()), "max": float(df[col].max()), "mean": float(df[col].mean())} for col in df.select_dtypes(include='number').columns[:10]},
    }
    prompt = f"""You are a senior data analyst. Answer this question about the dataset:

QUESTION: "{question}"
DATA INFO: {json.dumps(df_info, indent=2, default=str)}

RULES:
- Compute ACTUAL values from the data, not placeholders
- For customer-level metrics: deduplicate on customer_id if it exists
- For rate/percentage columns: use MEAN not SUM
- Assign the Plotly figure to 'result'

Respond with ONLY this JSON:
{{
  "answer": "Clear 2-3 sentence answer with actual numbers",
  "code": "Python using df, pd, px, go, np. Assign plotly figure to 'result'.",
  "insight": "One key business insight or recommendation"
}}"""
    try:
        response = client.messages.create(model="claude-sonnet-4-20250514", max_tokens=3000, messages=[{"role": "user", "content": prompt}])
        return extract_json(response.content[0].text)
    except Exception as e:
        st.error(f"API Error: {e}")
        return None


# ─────────────────────────────────────────────────────────
# DASHBOARD GENERATOR
# ─────────────────────────────────────────────────────────
def generate_dashboard(df, prompt, client):
    df_info = {
        "columns": list(df.columns), "dtypes": {col: str(dtype) for col, dtype in df.dtypes.items()},
        "shape": list(df.shape), "numeric_columns": list(df.select_dtypes(include=['number']).columns),
        "categorical_columns": list(df.select_dtypes(include=['object']).columns),
        "sample_values": {col: df[col].dropna().unique()[:5].tolist() for col in df.columns[:15]},
        "numeric_stats": {col: {"min": float(df[col].min()), "max": float(df[col].max()), "mean": round(float(df[col].mean()),2)} for col in df.select_dtypes(include='number').columns[:10]},
    }
    if 'customer_id' in df.columns:
        df_info['note'] = f"Dataset has {len(df)} rows but {df['customer_id'].nunique()} unique customers. Deduplicate on customer_id for customer-level KPIs."

    bq_context = ""
    bq = st.session_state.get('business_question')
    if bq:
        bq_context = f"\n\nBUSINESS QUESTION TO ANSWER: '{bq}'\nFocus all KPIs, charts, and insights on answering this specific question."

    ai_prompt = f"""You are an expert BI dashboard designer. Create a complete dashboard.

USER REQUEST: {prompt}{bq_context}

DATASET:
{json.dumps(df_info, indent=2, default=str)}

RESPOND WITH ONLY THIS JSON:
{{
  "title": "Dashboard Title",
  "description": "What this dashboard shows",
  "kpis": [{{"label": "KPI Name", "code": "result = df['col'].sum()", "format": "${{:,.0f}}"}}],
  "insights": ["Key insight with specific numbers"],
  "visualizations": [{{
    "title": "Chart Title",
    "chart_type": "bar",
    "code": "agg = df.groupby('col')['val'].mean().reset_index()\\nresult = px.bar(agg, x='col', y='val', color='col', color_discrete_sequence={VIVID_COLORS})",
    "insight": "What this reveals"
  }}]
}}

CRITICAL RULES:
1. Generate 4-6 visualizations with DIFFERENT chart types
2. Each visualization code must assign a Plotly figure to 'result'
3. Use ONLY actual column names from the dataset
4. KPI code must compute a real value and assign to 'result'
5. For customer-level KPIs: deduplicate first:
   customers = df.drop_duplicates(subset='customer_id', keep='last') if 'customer_id' in df.columns else df
6. For MAP: result = make_us_map(agg, value_col='revenue', state_col='state', title='Revenue by State', colorscale='Blues')
7. For RATES/PERCENTAGES: always use MEAN not SUM. Never sum percentages.
8. Cap displayed percentages at 100%.
9. COLOUR RULES — CRITICAL: Every chart must be vivid and multi-coloured:
   - Bar: px.bar(..., color='category_col', color_discrete_sequence=VIVID_COLORS)
   - Pie: px.pie(..., color_discrete_sequence=VIVID_COLORS)
   - Scatter: px.scatter(..., color='category_col', color_discrete_sequence=VIVID_COLORS)
   - Box: px.box(..., color='category_col', color_discrete_sequence=VIVID_COLORS)
   - Histogram: px.histogram(..., color_discrete_sequence=[VIVID_COLORS[2]])
   - Line: px.line(..., color='series_col', color_discrete_sequence=VIVID_COLORS)
   VIVID_COLORS = {VIVID_COLORS}
10. Return ONLY JSON"""

    try:
        response = client.messages.create(model="claude-sonnet-4-20250514", max_tokens=6000, messages=[{"role": "user", "content": ai_prompt}])
        result = extract_json(response.content[0].text)
        if result is None:
            st.error("⚠️ Could not parse dashboard JSON from AI response.")
            return None
        if 'visualizations' not in result:
            st.error("⚠️ AI response missing 'visualizations' key.")
            return None
        return result
    except Exception as e:
        st.error(f"API Error: {e}")
        return None


# ─────────────────────────────────────────────────────────
# RENDER DASHBOARD
# ─────────────────────────────────────────────────────────
def render_dashboard(df, dashboard):
    with st.sidebar:
        st.markdown("### Filters & Style")
        with st.expander("Chart Style", expanded=False):
            show_grid = st.checkbox("Show gridlines", value=False)
            height = st.slider("Chart height", 300, 700, 420, step=20)
            st.session_state['chart_style']['show_gridlines'] = show_grid
            st.session_state['chart_style']['chart_height'] = height
        st.markdown("---")
        st.markdown("**Data Filters**")
        slicer_cols = [c for c in df.columns if df[c].dtype == 'object' and 2 <= df[c].nunique() <= 20 and 'id' not in c.lower()][:6]
        filters = {}
        for col in slicer_cols:
            sel = st.multiselect(col.replace('_',' ').title(), sorted(df[col].dropna().unique().tolist()), key=f"slicer_{col}")
            if sel: filters[col] = sel
        st.session_state['filters'] = filters
        if filters: st.success(f"✅ {len(filters)} filter(s) active")
        if st.button("Clear Filters"): st.session_state['filters'] = {}; st.rerun()

    filtered_df = df.copy()
    for col, vals in st.session_state['filters'].items():
        if vals and col in filtered_df.columns:
            filtered_df = filtered_df[filtered_df[col].isin(vals)]
    if st.session_state['filters']:
        st.info(f"🔍 Showing **{len(filtered_df):,}** / {len(df):,} rows after filtering")

    st.markdown(f"## {dashboard.get('title', 'Dashboard')}")
    bq = st.session_state.get('business_question')
    if bq:
        st.markdown(f"""<div style='background:#0d0f1a;border:1px solid #1e2035;border-left:3px solid #6366f1;border-radius:6px;padding:0.5rem 0.9rem;margin-bottom:0.5rem'>
        <span style='color:#6366f1;font-size:0.7rem;font-weight:700;text-transform:uppercase;letter-spacing:0.08em'>Answering</span>
        <span style='color:#94a3b8;font-size:0.825rem;margin-left:0.5rem'>{bq}</span></div>""", unsafe_allow_html=True)
    if dashboard.get('description'):
        st.markdown(f"<span style='color:#64748b;font-size:0.875rem'>{dashboard.get('description','')}</span>", unsafe_allow_html=True)

    ec1, ec2, ec3, ec4 = st.columns(4)
    with ec1:
        if st.button("📄 Export PDF", use_container_width=True):
            with st.spinner("Building PDF..."):
                buf = export_to_pdf(dashboard, filtered_df)
            if buf:
                st.download_button("⬇️ Download PDF", buf, f"dashboard_{datetime.now().strftime('%Y%m%d_%H%M')}.pdf", "application/pdf", use_container_width=True)
    with ec2:
        if st.button("📊 Export PowerPoint", use_container_width=True):
            with st.spinner("Building PPTX..."):
                buf = export_to_powerpoint(dashboard, filtered_df)
            if buf:
                st.download_button("⬇️ Download PPTX", buf, f"dashboard_{datetime.now().strftime('%Y%m%d_%H%M')}.pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation", use_container_width=True)
    with ec3:
        if st.button("💾 Save Dashboard", use_container_width=True):
            st.session_state['saved_dashboards'].append({'name': dashboard.get('title','Dashboard'), 'timestamp': datetime.now().strftime('%Y-%m-%d %H:%M'), 'config': dashboard})
            st.success("✅ Saved!")
    with ec4:
        st.download_button("📋 JSON Config", json.dumps(dashboard, indent=2), f"dashboard_config_{datetime.now().strftime('%Y%m%d')}.json", "application/json", use_container_width=True)

    st.divider()
    kpis = dashboard.get('kpis', [])
    if kpis:
        st.markdown("### Key Metrics")
        cols = st.columns(len(kpis))
        for i, kpi in enumerate(kpis):
            with cols[i]:
                out = safe_exec(kpi.get('code', 'result=0'), filtered_df)
                if out['error']:
                    st.metric(kpi.get('label','KPI'), "⚠️ Error")
                else:
                    val = out['result']
                    try: display = kpi.get('format', '{:,}').format(val)
                    except: display = str(val)
                    st.metric(kpi.get('label','KPI'), display)
        st.divider()

    insights = dashboard.get('insights', [])
    if insights:
        st.markdown("### Key Insights")
        for ins in insights:
            st.markdown(f"""<div style='background:#12131f;border-left:2px solid #6366f1;border-radius:0 8px 8px 0;padding:0.6rem 1rem;margin:0.25rem 0;color:#94a3b8;font-size:0.875rem'>{ins}</div>""", unsafe_allow_html=True)
        st.divider()

    st.markdown("### Visualizations")
    vizs = dashboard.get('visualizations', [])
    chart_h = st.session_state['chart_style']['chart_height']
    for i in range(0, len(vizs), 2):
        cols = st.columns(2)
        for j in range(2):
            if i + j >= len(vizs): break
            viz = vizs[i + j]
            if not isinstance(viz, dict): continue
            with cols[j]:
                st.markdown(f"""<div style='background:#12131f;border:1px solid #1e2035;border-radius:12px;padding:0.75rem 1rem 0.25rem 1rem;margin-bottom:0.5rem'>
                <div style='font-size:0.75rem;font-weight:600;color:#64748b;text-transform:uppercase;letter-spacing:0.06em'>{viz.get('title', f'Chart {i+j+1}')}</div></div>""", unsafe_allow_html=True)
                out = safe_exec(viz.get('code', ''), filtered_df)
                if out['error']:
                    st.warning(f"⚠️ {out['error'][:120]}")
                    with st.expander("Debug"):
                        st.code(viz.get('code',''), language='python')
                        st.error(out['error'])
                elif out['result'] is not None and hasattr(out['result'], 'update_layout'):
                    fig = style_fig(out['result'], height=chart_h)
                    st.plotly_chart(fig, use_container_width=True, key=f"ch_{i}_{j}")
                elif isinstance(out['result'], pd.DataFrame):
                    st.dataframe(out['result'], use_container_width=True)
                if viz.get('insight'):
                    st.markdown(f"""<div style='background:#0d0e1a;border-left:2px solid #6366f1;border-radius:0 6px 6px 0;padding:0.35rem 0.75rem;margin-top:-0.25rem;margin-bottom:0.6rem'>
                    <span style='color:#94a3b8;font-size:0.75rem'>💡 {viz["insight"]}</span></div>""", unsafe_allow_html=True)


# ─────────────────────────────────────────────────────────
# PRESET DASHBOARDS
# ─────────────────────────────────────────────────────────
def generate_preset_dashboards(df, client):
    df_info = {
        "columns": list(df.columns), "dtypes": {col: str(dtype) for col, dtype in df.dtypes.items()},
        "numeric_columns": list(df.select_dtypes(include=['number']).columns),
        "categorical_columns": list(df.select_dtypes(include=['object']).columns),
        "shape": list(df.shape),
        "sample_values": {col: df[col].dropna().unique()[:5].tolist() for col in df.columns[:15]},
        "numeric_stats": {col: {"min": float(df[col].min()), "max": float(df[col].max()), "mean": round(float(df[col].mean()),2)} for col in df.select_dtypes(include='number').columns[:10]},
    }
    prompt = f"""Generate 5 role-based dashboard configs for this dataset.

DATASET: {json.dumps(df_info, indent=2, default=str)}

Respond ONLY with this JSON:
{{
  "dashboards": [
    {{
      "name": "Executive Overview", "icon": "👔", "audience": "C-Suite",
      "description": "High-level business health",
      "kpis": [{{"label": "KPI", "code": "result = df['col'].nunique()", "format": "{{:,}}"}}],
      "insights": ["insight with real numbers"],
      "visualizations": [{{
        "title": "Chart Title", "chart_type": "bar",
        "code": "agg = df.groupby('x')['y'].sum().reset_index(); result = px.bar(agg, x='x', y='y', color='x', color_discrete_sequence={VIVID_COLORS})",
        "insight": "What this shows"
      }}]
    }}
  ]
}}

5 dashboards: Executive Overview, Telecom Operations, Analytics Deep-Dive, Marketing & Growth, Financial Summary.

RULES:
- 3-4 KPIs per dashboard, 4-5 visualizations each, 2-3 insights
- For rates/percentages: ALWAYS use MEAN not SUM
- For map: result = make_us_map(agg, value_col='col', state_col='state', title='Title', colorscale='Blues')
- COLOUR RULES — all charts must be vivid multi-coloured:
  * Bar: px.bar(..., color='category_col', color_discrete_sequence=VIVID_COLORS)
  * Pie: px.pie(..., color_discrete_sequence=VIVID_COLORS)
  * Scatter: px.scatter(..., color='category_col', color_discrete_sequence=VIVID_COLORS)
  * Box: px.box(..., color='category_col', color_discrete_sequence=VIVID_COLORS)
  * Histogram: px.histogram(..., color_discrete_sequence=[VIVID_COLORS[2]])
  VIVID_COLORS = {VIVID_COLORS}
- Assign all figures to 'result'
- Return ONLY JSON"""
    try:
        response = client.messages.create(model="claude-sonnet-4-20250514", max_tokens=8000, messages=[{"role": "user", "content": prompt}])
        result = extract_json(response.content[0].text)
        if result and 'dashboards' in result:
            return result
        st.error("⚠️ Could not parse preset dashboards.")
        return None
    except Exception as e:
        st.error(f"API Error: {e}")
        return None


# ─────────────────────────────────────────────────────────
# RECOMMENDATIONS ENGINE
# ─────────────────────────────────────────────────────────
PRIORITY_COLORS = {"critical": "#ef4444", "high": "#f97316", "medium": "#f59e0b", "low": "#10b981"}

def run_recommendations_engine(df, client, analyst_report=None, business_question=None):
    df_info = {
        "columns": list(df.columns), "shape": list(df.shape),
        "numeric_stats": {col: {"min": float(df[col].min()), "max": float(df[col].max()), "mean": round(float(df[col].mean()),2)} for col in df.select_dtypes(include='number').columns[:10]},
        "sample_values": {col: df[col].dropna().unique()[:5].tolist() for col in df.columns[:10]},
    }
    prior_insights = []
    if analyst_report:
        prior_insights = [i.get('insight','') for i in analyst_report.get('analysis_insights', [])]

    prompt = f"""You are a Chief Data Officer. Based on the data analysis below, provide strategic business recommendations.

BUSINESS QUESTION: {business_question or 'General business optimization'}
DATA PROFILE: {json.dumps(df_info, indent=2, default=str)}
PRIOR INSIGHTS: {json.dumps(prior_insights, indent=2)}

Respond ONLY with this JSON:
{{
  "executive_brief": "2-3 sentence executive summary of the situation and key opportunity",
  "summary_stats": {{
    "total_recommendations": 6,
    "quick_wins": 2,
    "strategic_initiatives": 2,
    "risks_identified": 2,
    "estimated_impact": "High"
  }},
  "recommendations": [
    {{
      "id": 1,
      "title": "Recommendation title",
      "category": "Revenue / Customer / Operations / Risk / Growth / Cost",
      "priority": "critical / high / medium / low",
      "effort": "low / medium / high",
      "impact": "low / medium / high",
      "description": "Detailed description of what to do and why",
      "action_items": ["Specific action 1", "Specific action 2", "Specific action 3"],
      "expected_outcome": "Quantified expected outcome",
      "timeframe": "1-2 weeks / 1-3 months / 3-6 months / 6-12 months",
      "auto_applicable": true
    }}
  ],
  "quick_wins": ["Quick win 1 — can be automated", "Quick win 2"],
  "risks_if_ignored": [
    {{"risk": "What happens if not addressed", "severity": "high / medium / low"}}
  ],
  "success_metrics": ["Metric 1 to track", "Metric 2", "Metric 3"]
}}

Generate 5-7 recommendations. Mark auto_applicable=true only for data cleaning/transformation steps."""
    try:
        response = client.messages.create(model="claude-sonnet-4-20250514", max_tokens=5000, messages=[{"role": "user", "content": prompt}])
        return extract_json(response.content[0].text)
    except Exception as e:
        st.error(f"Recommendations error: {e}")
        return None


def render_recommendations(rec, df, client):
    if not rec:
        return

    exec_brief = rec.get("executive_brief", "")
    if exec_brief:
        st.markdown(f"""<div style='background:#12131f;border:1px solid #1e2035;border-left:3px solid #6366f1;border-radius:8px;padding:1rem 1.25rem;margin:0.5rem 0'>
        <div style='font-size:0.68rem;font-weight:700;color:#6366f1;text-transform:uppercase;letter-spacing:0.1em;margin-bottom:0.4rem'>Executive Brief</div>
        <div style='color:#94a3b8;font-size:0.875rem;line-height:1.5'>{exec_brief}</div></div>""", unsafe_allow_html=True)

    ss = rec.get("summary_stats", {})
    if ss:
        s1, s2, s3, s4, s5 = st.columns(5)
        s1.metric("Total Recs", ss.get("total_recommendations", 0))
        s2.metric("Quick Wins", ss.get("quick_wins", 0))
        s3.metric("Strategic", ss.get("strategic_initiatives", 0))
        s4.metric("Risks", ss.get("risks_identified", 0))
        s5.metric("Est. Impact", ss.get("estimated_impact", "—"))

    st.divider()

    recs = rec.get("recommendations", [])
    if recs:
        effort_map = {"low": 1, "medium": 2, "high": 3}
        impact_map = {"low": 1, "medium": 2, "high": 3}
        scatter_data = []
        for r in recs:
            scatter_data.append({
                "title": r.get("title","")[:30], "effort": effort_map.get(r.get("effort","medium"), 2),
                "impact": impact_map.get(r.get("impact","medium"), 2),
                "priority": r.get("priority","medium"), "category": r.get("category","General"),
            })
        sdf = pd.DataFrame(scatter_data)
        fig = px.scatter(sdf, x="effort", y="impact", color="priority", text="title", size=[20]*len(sdf),
                        color_discrete_map={"critical":"#ef4444","high":"#f97316","medium":"#f59e0b","low":"#10b981"},
                        labels={"effort":"Effort","impact":"Impact"}, title="Effort × Impact Matrix")
        fig.update_traces(textposition='top center', textfont=dict(size=10, color='#94a3b8'))
        fig.update_xaxes(tickvals=[1,2,3], ticktext=["Low","Medium","High"])
        fig.update_yaxes(tickvals=[1,2,3], ticktext=["Low","Medium","High"])
        fig = style_fig(fig, height=380)
        st.plotly_chart(fig, use_container_width=True)
        st.divider()

        cat_colors = {"Revenue":"#6366f1","Customer":"#8b5cf6","Operations":"#06b6d4","Risk":"#ef4444","Growth":"#10b981","Cost":"#f59e0b"}
        by_priority = {}
        for r in recs:
            p = r.get("priority","medium")
            by_priority.setdefault(p, []).append(r)
        for priority in ["critical","high","medium","low"]:
            group = by_priority.get(priority, [])
            if not group: continue
            p_color = PRIORITY_COLORS.get(priority, "#64748b")
            st.markdown(f"<div style='font-size:0.7rem;font-weight:700;color:{p_color};text-transform:uppercase;letter-spacing:0.1em;margin:0.75rem 0 0.35rem 0'>{priority.upper()} PRIORITY</div>", unsafe_allow_html=True)
            for r in group:
                cat = r.get("category","General")
                cat_color = cat_colors.get(cat, "#64748b")
                actions_html = "".join([f"<div style='font-size:0.78rem;color:#64748b;padding:0.15rem 0'>• {a}</div>" for a in r.get("action_items",[])])
                st.markdown(f"""<div style='background:#12131f;border:1px solid #1e2035;border-radius:10px;padding:1rem 1.2rem;margin:0.35rem 0'>
                <div style='display:flex;justify-content:space-between;align-items:flex-start;margin-bottom:0.5rem'>
                  <div>
                    <span style='color:{cat_color};font-size:0.65rem;font-weight:700;text-transform:uppercase;letter-spacing:0.08em'>{cat}</span>
                    <div style='color:#e2e8f0;font-weight:700;font-size:0.9rem;margin-top:0.15rem'>{r.get("title","")}</div>
                  </div>
                  <div style='text-align:right;flex-shrink:0;margin-left:1rem'>
                    <span style='background:{p_color}20;color:{p_color};font-size:0.65rem;font-weight:700;padding:0.2rem 0.5rem;border-radius:4px;text-transform:uppercase'>{priority}</span>
                    <div style='color:#475569;font-size:0.7rem;margin-top:0.25rem'>{r.get("timeframe","")}</div>
                  </div>
                </div>
                <div style='color:#64748b;font-size:0.825rem;margin-bottom:0.5rem'>{r.get("description","")}</div>
                {actions_html}
                <div style='margin-top:0.5rem;padding-top:0.5rem;border-top:1px solid #1e2035'>
                  <span style='color:#10b981;font-size:0.78rem'>→ {r.get("expected_outcome","")}</span>
                </div>
                </div>""", unsafe_allow_html=True)

    st.divider()

    quick_wins = rec.get("quick_wins", [])
    if quick_wins:
        st.markdown("#### ⚡ Quick Wins")
        for i, qw in enumerate(quick_wins):
            st.markdown(f"<div style='background:#12131f;border:1px solid #1e2035;border-radius:7px;padding:0.5rem 0.9rem;margin:0.2rem 0;font-size:0.825rem;color:#94a3b8'><span style='font-weight:700;color:#10b981'>{i+1}.</span> {qw}</div>", unsafe_allow_html=True)

        if st.button("⚡ Auto-Apply All Quick Wins", type="primary", use_container_width=True, key="auto_apply_qw"):
            _df = st.session_state.get('cleaned_df') if st.session_state.get('cleaned_df') is not None else st.session_state.get('uploaded_df')
            if _df is None:
                st.warning("No data loaded.")
            else:
                with st.spinner("Applying quick wins..."):
                    _before_df = _df.copy()
                    _score_before = compute_quality_score(_before_df)
                    qw_prompt = f"""Apply these quick wins to the dataframe:
{json.dumps(quick_wins, indent=2)}

Dataset columns: {list(_df.columns)}
Dataset shape: {_df.shape}

Write Python code that:
1. Starts with: df_qw = df.copy()
2. Applies each quick win that is actually automatable (skip business strategy items)
3. Assigns result to: result = df_qw
4. Returns a JSON with keys:
   - "code": the Python cleaning code
   - "applied": list of quick wins that were automated
   - "skipped": list that require manual action

Respond ONLY with JSON."""
                    try:
                        _resp = client.messages.create(model="claude-sonnet-4-20250514", max_tokens=3000, messages=[{"role": "user", "content": qw_prompt}])
                        _parsed = extract_json(_resp.content[0].text)
                        if _parsed and _parsed.get('code'):
                            _out = _safe_clean_exec(_parsed['code'], _df)
                            if _out['error']:
                                st.error(f"Auto-apply error: {_out['error']}")
                            elif _out['result'] is not None and isinstance(_out['result'], pd.DataFrame):
                                _new_df = _out['result']
                                st.session_state['cleaned_df'] = _new_df
                                _score_after = compute_quality_score(_new_df)
                                st.session_state['quick_wins_result'] = {
                                    'applied': _parsed.get('applied', []),
                                    'skipped': _parsed.get('skipped', []),
                                    'score_before': _score_before,
                                    'score_after': _score_after,
                                    'rows_before': len(_before_df),
                                    'rows_after': len(_new_df),
                                    'cols_before': len(_before_df.columns),
                                    'cols_after': len(_new_df.columns),
                                }
                                st.rerun()
                    except Exception as ex:
                        st.error(f"Error: {ex}")

        qwr = st.session_state.get('quick_wins_result')
        if qwr:
            sb = qwr['score_before']; sa = qwr['score_after']; delta = sa - sb
            dc = "#10b981" if delta >= 0 else "#ef4444"; di = "↑" if delta >= 0 else "↓"
            applied = qwr.get('applied', []); skipped = qwr.get('skipped', [])
            st.markdown(f"""<div style='background:#12131f;border:1px solid #166534;border-radius:14px;padding:1.1rem 1.4rem;margin:0.75rem 0'>
              <div style='font-size:0.68rem;font-weight:700;color:#10b981;text-transform:uppercase;letter-spacing:0.1em;margin-bottom:0.75rem'>✅ Auto-Apply Results</div>
              <div style='display:grid;grid-template-columns:repeat(4,1fr);gap:1rem;margin-bottom:0.75rem'>
                <div style='text-align:center'><div style='font-size:0.65rem;color:#64748b;font-weight:600;text-transform:uppercase;margin-bottom:2px'>Quality Before</div><div style='font-size:1.8rem;font-weight:800;color:{"#10b981" if sb>=80 else "#f59e0b" if sb>=60 else "#ef4444"}'>{sb}%</div></div>
                <div style='text-align:center'><div style='font-size:0.65rem;color:#64748b;font-weight:600;text-transform:uppercase;margin-bottom:2px'>Quality After</div><div style='font-size:1.8rem;font-weight:800;color:{"#10b981" if sa>=80 else "#f59e0b" if sa>=60 else "#ef4444"}'>{sa}%</div></div>
                <div style='text-align:center'><div style='font-size:0.65rem;color:#64748b;font-weight:600;text-transform:uppercase;margin-bottom:2px'>Score Delta</div><div style='font-size:1.8rem;font-weight:800;color:{dc}'>{di}{abs(delta)} pts</div></div>
                <div style='text-align:center'><div style='font-size:0.65rem;color:#64748b;font-weight:600;text-transform:uppercase;margin-bottom:2px'>Rows</div><div style='font-size:1.2rem;font-weight:700;color:#e2e8f0'>{qwr['rows_before']:,} → {qwr['rows_after']:,}</div><div style='font-size:0.65rem;color:#64748b'>Cols: {qwr['cols_before']} → {qwr['cols_after']}</div></div>
              </div></div>""", unsafe_allow_html=True)
            if applied:
                st.markdown("<div style='margin:0.5rem 0 0.25rem 0;font-size:0.75rem;color:#10b981;font-weight:700;text-transform:uppercase;letter-spacing:0.06em'>✅ Applied Automatically</div>", unsafe_allow_html=True)
                for item in applied:
                    st.markdown(f"<div style='background:#0d2318;border:1px solid #166534;border-radius:6px;padding:0.4rem 0.8rem;margin:0.2rem 0;font-size:0.8rem;color:#86efac'>✔️ {item}</div>", unsafe_allow_html=True)
            if skipped:
                st.markdown("<div style='margin:0.5rem 0 0.25rem 0;font-size:0.75rem;color:#f59e0b;font-weight:700;text-transform:uppercase;letter-spacing:0.06em'>⚠️ Requires Manual Action</div>", unsafe_allow_html=True)
                for item in skipped:
                    st.markdown(f"<div style='background:#1f1707;border:1px solid #713f12;border-radius:6px;padding:0.4rem 0.8rem;margin:0.2rem 0;font-size:0.8rem;color:#fde68a'>👉 {item}</div>", unsafe_allow_html=True)
            if st.button("✕ Clear Result", key="clear_qwr"):
                st.session_state['quick_wins_result'] = None
                st.rerun()

    risks = rec.get("risks_if_ignored", [])
    if risks:
        st.divider()
        st.markdown("#### 🚨 Risks If Not Addressed")
        risk_cols = st.columns(min(len(risks), 3))
        for i, risk in enumerate(risks):
            sev = risk.get("severity","medium"); s_col = PRIORITY_COLORS.get(sev, "#f59e0b")
            with risk_cols[i % 3]:
                st.markdown(f"<div style='background:#12131f;border:1px solid #1e2035;border-top:2px solid {s_col};border-radius:9px;padding:0.75rem 0.9rem;margin-bottom:0.4rem'><div style='font-size:0.62rem;color:{s_col};font-weight:700;text-transform:uppercase;letter-spacing:0.06em;margin-bottom:0.3rem'>{sev.upper()} RISK</div><div style='font-size:0.78rem;color:#94a3b8;line-height:1.4'>{risk.get('risk','')}</div></div>", unsafe_allow_html=True)

    success_metrics = rec.get("success_metrics", [])
    if success_metrics:
        st.divider()
        st.markdown("#### 📏 Success Metrics to Track")
        sm_cols = st.columns(min(len(success_metrics), 5))
        for i, sm in enumerate(success_metrics):
            with sm_cols[i % 5]:
                st.markdown(f"<div style='background:#12131f;border:1px solid #1e2035;border-radius:8px;padding:0.6rem 0.8rem;text-align:center;font-size:0.75rem;color:#a5b4fc;font-weight:600;line-height:1.4'>📊 {sm}</div>", unsafe_allow_html=True)


# ─────────────────────────────────────────────────────────
# STEP 10 — MONITOR & ITERATE
# ─────────────────────────────────────────────────────────
def run_monitor_analysis(df, kpi_configs, snapshots, client, business_question=None):
    current_stats = {}
    for kpi in kpi_configs:
        try:
            out = safe_exec(kpi['code'], df)
            current_stats[kpi['label']] = round(float(out['result']), 2) if out['result'] is not None and not out['error'] else None
        except Exception:
            current_stats[kpi['label']] = None

    snap_summary = []
    for snap in snapshots[-3:]:
        snap_summary.append({'timestamp': snap['timestamp'], 'kpis': snap.get('kpis', {}), 'rows': snap.get('rows', 0), 'quality': snap.get('quality', 0)})

    prompt = f"""You are a data monitoring AI. Compare current metrics against historical snapshots and identify anomalies, trends, and recommended actions.

BUSINESS QUESTION: {business_question or 'General monitoring'}
CURRENT METRICS: {json.dumps(current_stats, indent=2)}
HISTORICAL SNAPSHOTS: {json.dumps(snap_summary, indent=2)}
CURRENT ROWS: {len(df)}
CURRENT QUALITY SCORE: {compute_quality_score(df)}

Respond ONLY with this JSON:
{{
  "status": "healthy / warning / critical",
  "headline": "One sentence summary of current data health",
  "anomalies": [
    {{"metric": "KPI name", "current": 0, "previous": 0, "change_pct": 0, "severity": "high/medium/low", "explanation": "Why this is notable"}}
  ],
  "trends": [
    {{"metric": "KPI name", "direction": "up/down/stable", "description": "Trend description"}}
  ],
  "actions": [
    {{"action": "Specific action to take", "priority": "high/medium/low", "timeframe": "immediate/this week/this month"}}
  ],
  "next_check_recommendation": "When to check again and what to watch"
}}"""
    try:
        response = client.messages.create(model="claude-sonnet-4-20250514", max_tokens=3000, messages=[{"role": "user", "content": prompt}])
        return extract_json(response.content[0].text)
    except Exception as e:
        return None


def render_monitor_tab(df, client):
    st.markdown('<div class="step-badge">Step 10 — Monitor & Iterate</div>', unsafe_allow_html=True)
    st.header("Monitor & Iterate")
    st.caption("Track KPI changes over time, detect anomalies, and get AI-driven iteration recommendations.")

    st.markdown("### ⚙️ Configure KPIs to Monitor")
    num_cols = list(df.select_dtypes(include='number').columns)

    with st.expander("Add / Edit Monitored KPIs", expanded=not bool(st.session_state['monitor_kpi_config'])):
        kpi_presets = []
        for col in num_cols[:6]:
            kpi_presets.append({'label': col.replace('_', ' ').title(), 'code': f"result = df['{col}'].mean()", 'format': '{:,.2f}'})
        if kpi_presets and not st.session_state['monitor_kpi_config']:
            if st.button("⚡ Auto-detect KPIs from data", type="primary"):
                st.session_state['monitor_kpi_config'] = kpi_presets
                st.rerun()

        st.markdown("**Manual KPI Setup:**")
        mc1, mc2, mc3 = st.columns([2, 3, 1])
        with mc1:
            new_label = st.text_input("KPI Label", placeholder="e.g., Avg Revenue", key="mon_kpi_label")
        with mc2:
            new_code  = st.text_input("Python code (assign to result)", placeholder="result = df['revenue'].mean()", key="mon_kpi_code")
        with mc3:
            st.markdown("<br>", unsafe_allow_html=True)
            if st.button("+ Add KPI", use_container_width=True):
                if new_label and new_code:
                    st.session_state['monitor_kpi_config'].append({'label': new_label, 'code': new_code, 'format': '{:,.2f}'})
                    st.rerun()

        if st.session_state['monitor_kpi_config']:
            st.markdown("**Configured KPIs:**")
            for i, kpi in enumerate(st.session_state['monitor_kpi_config']):
                kc1, kc2 = st.columns([4, 1])
                with kc1:
                    st.markdown(f"<span style='color:#a5b4fc;font-size:0.825rem;font-weight:600'>{kpi['label']}</span> <code style='font-size:0.75rem;color:#64748b'>{kpi['code']}</code>", unsafe_allow_html=True)
                with kc2:
                    if st.button("✕", key=f"del_kpi_{i}", use_container_width=True):
                        st.session_state['monitor_kpi_config'].pop(i)
                        st.rerun()

    st.markdown("### 📸 Snapshots")
    sc1, sc2, sc3 = st.columns([2, 2, 1])
    with sc1:
        snap_label = st.text_input("Snapshot label", placeholder="e.g., After Q2 cleaning", key="snap_label_input", label_visibility="collapsed")
    with sc2:
        if st.button("📸 Take Snapshot Now", type="primary", use_container_width=True):
            kpi_values = {}
            for kpi in st.session_state['monitor_kpi_config']:
                out = safe_exec(kpi['code'], df)
                kpi_values[kpi['label']] = round(float(out['result']), 4) if out['result'] is not None and not out['error'] else None
            snapshot = {
                'timestamp': datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
                'label': snap_label or f"Snapshot {len(st.session_state['monitor_snapshots'])+1}",
                'rows': len(df), 'columns': len(df.columns),
                'quality': compute_quality_score(df),
                'kpis': kpi_values,
            }
            st.session_state['monitor_snapshots'].append(snapshot)
            st.success(f"✅ Snapshot '{snapshot['label']}' saved!")
            st.rerun()
    with sc3:
        if st.session_state['monitor_snapshots']:
            if st.button("Clear All", use_container_width=True, key="clear_snapshots"):
                st.session_state['monitor_snapshots'] = []
                st.rerun()

    snapshots = st.session_state['monitor_snapshots']
    if snapshots:
        st.markdown(f"**{len(snapshots)} snapshot(s) recorded**")
        kpi_config = st.session_state['monitor_kpi_config']
        if kpi_config:
            st.markdown("### 📊 Live KPI Dashboard")
            live_cols = st.columns(min(len(kpi_config), 4))
            for i, kpi in enumerate(kpi_config[:4]):
                with live_cols[i]:
                    out = safe_exec(kpi['code'], df)
                    current_val = round(float(out['result']), 2) if out['result'] is not None and not out['error'] else None
                    prev_val = None
                    if snapshots:
                        prev_val = snapshots[-1]['kpis'].get(kpi['label'])
                    delta_display = None
                    if current_val is not None and prev_val is not None:
                        delta_display = round(current_val - prev_val, 2)
                    try:
                        display = kpi.get('format', '{:,.2f}').format(current_val) if current_val is not None else 'N/A'
                    except Exception:
                        display = str(current_val)
                    st.metric(kpi['label'], display, delta=f"{delta_display:+,.2f}" if delta_display is not None else None)

        snap_data = []
        for snap in snapshots:
            row = {'Snapshot': snap.get('label',''), 'Timestamp': snap['timestamp'], 'Rows': f"{snap['rows']:,}", 'Cols': snap['columns'], 'Quality': f"{snap['quality']}%"}
            for k, v in snap.get('kpis', {}).items():
                row[k] = round(v, 2) if v is not None else 'N/A'
            snap_data.append(row)
        snap_df = pd.DataFrame(snap_data)
        st.dataframe(snap_df, use_container_width=True)

        if len(snapshots) >= 2:
            st.markdown("### 📈 Quality Score Over Time")
            trend_df = pd.DataFrame([{'Snapshot': s.get('label', f"#{i+1}"), 'Quality Score': s['quality'], 'Rows': s['rows']} for i, s in enumerate(snapshots)])
            fig_trend = px.line(trend_df, x='Snapshot', y='Quality Score', markers=True, title='Data Quality Score Over Snapshots', color_discrete_sequence=[VIVID_COLORS[0]])
            fig_trend.add_hline(y=80, line_dash='dash', line_color='#10b981', annotation_text='Good (80%)', annotation_font_color='#10b981')
            fig_trend.add_hline(y=60, line_dash='dash', line_color='#f59e0b', annotation_text='Acceptable (60%)', annotation_font_color='#f59e0b')
            fig_trend = style_fig(fig_trend, height=300)
            st.plotly_chart(fig_trend, use_container_width=True)

            if kpi_config:
                st.markdown("### 📊 KPI Trend Over Snapshots")
                kpi_trend_rows = []
                for snap in snapshots:
                    for kpi_name, val in snap.get('kpis', {}).items():
                        if val is not None:
                            kpi_trend_rows.append({'Snapshot': snap.get('label', snap['timestamp']), 'KPI': kpi_name, 'Value': val})
                if kpi_trend_rows:
                    kpi_trend_df = pd.DataFrame(kpi_trend_rows)
                    valid_kpis = kpi_trend_df.groupby('KPI').size()
                    valid_kpis = valid_kpis[valid_kpis >= 2].index.tolist()
                    if valid_kpis:
                        kpi_trend_df = kpi_trend_df[kpi_trend_df['KPI'].isin(valid_kpis)]
                        fig_kpi = px.line(kpi_trend_df, x='Snapshot', y='Value', color='KPI',
                                         title='KPI Values Over Snapshots', markers=True,
                                         color_discrete_sequence=VIVID_COLORS)
                        fig_kpi = style_fig(fig_kpi, height=360)
                        st.plotly_chart(fig_kpi, use_container_width=True)

        st.divider()
        st.markdown("### 🤖 AI Monitoring Analysis")
        if st.session_state['monitor_kpi_config']:
            if st.button("🔍 Run AI Monitoring Analysis", type="primary", use_container_width=True, key="run_monitor_ai"):
                with st.spinner("AI analyzing trends and anomalies..."):
                    mon_result = run_monitor_analysis(
                        df, st.session_state['monitor_kpi_config'], snapshots, client,
                        business_question=st.session_state.get('business_question')
                    )
                    if mon_result:
                        st.session_state['monitor_ai_result'] = mon_result
                        st.rerun()

            mon_ai = st.session_state.get('monitor_ai_result')
            if mon_ai:
                status = mon_ai.get('status', 'unknown')
                status_colors = {'healthy': '#10b981', 'warning': '#f59e0b', 'critical': '#ef4444'}
                status_icons  = {'healthy': '✅', 'warning': '⚠️', 'critical': '🚨'}
                sc = status_colors.get(status, '#64748b')
                si = status_icons.get(status, 'ℹ️')
                st.markdown(f"""
<div style='background:#12131f;border:1px solid {sc};border-radius:12px;padding:1rem 1.3rem;margin:0.5rem 0'>
  <div style='display:flex;align-items:center;gap:0.75rem'>
    <div style='font-size:1.5rem'>{si}</div>
    <div>
      <div style='font-size:0.65rem;font-weight:700;color:{sc};text-transform:uppercase;letter-spacing:0.1em'>{status.upper()}</div>
      <div style='font-size:0.9rem;color:#e2e8f0;font-weight:600;margin-top:0.1rem'>{mon_ai.get('headline','')}</div>
    </div>
  </div>
</div>""", unsafe_allow_html=True)

                anomalies = mon_ai.get('anomalies', [])
                if anomalies:
                    st.markdown("#### 🔴 Anomalies Detected")
                    for anom in anomalies:
                        sev = anom.get('severity', 'medium')
                        a_color = PRIORITY_COLORS.get(sev, '#64748b')
                        chg = anom.get('change_pct', 0)
                        chg_color = '#10b981' if chg >= 0 else '#ef4444'
                        st.markdown(f"""
<div style='background:#12131f;border:1px solid #1e2035;border-left:2px solid {a_color};border-radius:8px;padding:0.75rem 1rem;margin:0.3rem 0'>
  <div style='display:flex;justify-content:space-between;align-items:flex-start'>
    <div>
      <span style='color:{a_color};font-size:0.65rem;font-weight:700;text-transform:uppercase'>{sev}</span>
      <div style='color:#e2e8f0;font-weight:600;font-size:0.875rem'>{anom.get('metric','')}</div>
      <div style='color:#64748b;font-size:0.78rem;margin-top:0.25rem'>{anom.get('explanation','')}</div>
    </div>
    <div style='text-align:right;flex-shrink:0;margin-left:1rem'>
      <div style='color:#94a3b8;font-size:0.75rem'>Prev: <strong>{anom.get('previous','—')}</strong></div>
      <div style='color:#94a3b8;font-size:0.75rem'>Now: <strong>{anom.get('current','—')}</strong></div>
      <div style='color:{chg_color};font-size:0.78rem;font-weight:700'>{'+' if chg >= 0 else ''}{chg:.1f}%</div>
    </div>
  </div>
</div>""", unsafe_allow_html=True)

                trends = mon_ai.get('trends', [])
                if trends:
                    st.markdown("#### 📈 Trend Analysis")
                    trend_cols = st.columns(min(len(trends), 3))
                    dir_icons = {'up': '↑', 'down': '↓', 'stable': '→'}
                    dir_colors = {'up': '#10b981', 'down': '#ef4444', 'stable': '#6366f1'}
                    for i, tr in enumerate(trends):
                        d = tr.get('direction', 'stable')
                        with trend_cols[i % 3]:
                            st.markdown(f"""
<div style='background:#12131f;border:1px solid #1e2035;border-radius:9px;padding:0.75rem 0.9rem;margin-bottom:0.4rem'>
  <div style='display:flex;align-items:center;gap:0.5rem;margin-bottom:0.25rem'>
    <span style='font-size:1.1rem;color:{dir_colors.get(d,"#6366f1")}'>{dir_icons.get(d,"→")}</span>
    <span style='color:#e2e8f0;font-weight:600;font-size:0.825rem'>{tr.get('metric','')}</span>
  </div>
  <div style='color:#64748b;font-size:0.75rem'>{tr.get('description','')}</div>
</div>""", unsafe_allow_html=True)

                actions = mon_ai.get('actions', [])
                if actions:
                    st.markdown("#### 🎯 Recommended Actions")
                    for act in actions:
                        p = act.get('priority', 'medium')
                        p_color = PRIORITY_COLORS.get(p, '#64748b')
                        st.markdown(f"""
<div style='background:#12131f;border:1px solid #1e2035;border-radius:7px;padding:0.6rem 0.9rem;margin:0.25rem 0;display:flex;justify-content:space-between;align-items:center'>
  <span style='color:#94a3b8;font-size:0.825rem'>{act.get('action','')}</span>
  <div style='flex-shrink:0;margin-left:0.75rem;text-align:right'>
    <span style='background:{p_color}20;color:{p_color};font-size:0.62rem;font-weight:700;padding:0.15rem 0.45rem;border-radius:4px;text-transform:uppercase'>{p}</span>
    <div style='color:#475569;font-size:0.65rem;margin-top:0.15rem'>{act.get('timeframe','')}</div>
  </div>
</div>""", unsafe_allow_html=True)

                nxt = mon_ai.get('next_check_recommendation', '')
                if nxt:
                    st.info(f"🔔 **Next Check:** {nxt}")

                if st.button("✕ Clear AI Analysis", key="clear_mon_ai"):
                    st.session_state['monitor_ai_result'] = None
                    st.rerun()
        else:
            st.info("👆 Configure KPIs above and take at least one snapshot before running AI analysis.")

    else:
        st.markdown("""
<div style='background:#12131f;border:1px dashed #2d3050;border-radius:12px;padding:2rem;text-align:center;margin:1rem 0'>
  <div style='font-size:2rem;margin-bottom:0.75rem'>📸</div>
  <div style='color:#e2e8f0;font-weight:600;font-size:1rem;margin-bottom:0.4rem'>No snapshots yet</div>
  <div style='color:#64748b;font-size:0.825rem'>Configure KPIs above, then take a snapshot to start monitoring.</div>
</div>""", unsafe_allow_html=True)

    st.divider()
    st.markdown("### 📓 Iteration Log")
    iter_note = st.text_area("Log a note", placeholder="e.g., Re-ran cleaning after fixing source data.", height=80, key="iter_note_input", label_visibility="collapsed")
    ic1, ic2 = st.columns([3, 1])
    with ic1:
        if st.button("📝 Save Note", use_container_width=True, key="save_iter_note"):
            if iter_note.strip():
                if 'iteration_log' not in st.session_state:
                    st.session_state['iteration_log'] = []
                st.session_state['iteration_log'].append({
                    'timestamp': datetime.now().strftime('%Y-%m-%d %H:%M'),
                    'note': iter_note.strip(),
                    'data_rows': len(df),
                    'quality': compute_quality_score(df),
                })
                st.rerun()
    with ic2:
        if 'iteration_log' in st.session_state and st.session_state['iteration_log']:
            if st.button("Clear Log", use_container_width=True, key="clear_iter_log"):
                st.session_state['iteration_log'] = []
                st.rerun()

    iter_log = st.session_state.get('iteration_log', [])
    if iter_log:
        for entry in reversed(iter_log):
            st.markdown(f"""
<div style='background:#12131f;border:1px solid #1e2035;border-radius:8px;padding:0.6rem 0.9rem;margin:0.2rem 0;display:flex;gap:1rem;align-items:flex-start'>
  <div style='flex-shrink:0'>
    <div style='color:#6366f1;font-size:0.65rem;font-weight:700'>{entry['timestamp']}</div>
    <div style='color:#475569;font-size:0.62rem'>{entry['data_rows']:,} rows · Q:{entry['quality']}%</div>
  </div>
  <div style='color:#94a3b8;font-size:0.8rem;line-height:1.45'>{entry['note']}</div>
</div>""", unsafe_allow_html=True)


# ─────────────────────────────────────────────────────────
# MAIN — TAB LAYOUT
# ─────────────────────────────────────────────────────────
def _render_step1_bq():
    st.markdown('<div class="step-badge">Step 1 — Define Your Question</div>', unsafe_allow_html=True)
    st.header("What are you trying to figure out?")
    st.markdown("<span style='color:#64748b;font-size:0.875rem'>Start with a business question — DashAI will tailor every analysis and dashboard to answer it.</span>", unsafe_allow_html=True)

    example_qs = ["Why is customer churn increasing?", "Which products are most profitable by region?", "What's driving revenue growth this quarter?", "Which customer segments have the highest lifetime value?", "Where are we losing customers in the funnel?", "Which states/regions have the best performance?"]
    eq_cols = st.columns(3)
    for i, eq in enumerate(example_qs):
        with eq_cols[i % 3]:
            if st.button(eq, key=f"eq_{i}", use_container_width=True):
                st.session_state['business_question'] = eq
                st.session_state['bq_confirmed'] = False
                st.rerun()

    bq_input = st.text_area("Your business question", value=st.session_state.get('business_question') or '', placeholder="e.g., Why are high-value customers churning after month 3?", height=80, label_visibility="collapsed")

    bq_col1, bq_col2 = st.columns([3, 1])
    with bq_col1:
        if st.button("✓ Set Business Question", type="primary", use_container_width=True):
            if bq_input.strip():
                st.session_state['business_question'] = bq_input.strip()
                st.session_state['bq_confirmed'] = True
                st.session_state['analyst_report'] = None
                st.session_state['dashboard_result'] = None
                st.session_state['prompt_template'] = bq_input.strip()
                st.rerun()
            else:
                st.warning("Please enter or pick a business question.")
    with bq_col2:
        if st.session_state.get('business_question'):
            if st.button("Clear", key="clear_bq", use_container_width=True):
                st.session_state['business_question'] = None
                st.session_state['bq_confirmed'] = False
                st.rerun()

    if st.session_state.get('bq_confirmed') and st.session_state.get('business_question'):
        bq = st.session_state['business_question']
        st.markdown(f"""<div style='background:#0d1a0f;border:1px solid #166534;border-left:3px solid #10b981;border-radius:8px;padding:0.75rem 1rem;margin:0.75rem 0'>
        <span style='color:#10b981;font-size:0.7rem;font-weight:700;text-transform:uppercase;letter-spacing:0.08em'>✓ Business Question Set</span><br>
        <span style='color:#e2e8f0;font-size:0.9rem;font-weight:500'>{bq}</span></div>""", unsafe_allow_html=True)
    elif st.session_state.get('business_question') and not st.session_state.get('bq_confirmed'):
        st.info("👆 Click **Set Business Question** to confirm and continue")
    else:
        st.markdown("""<div style='background:#12131f;border:1px dashed #2d3050;border-radius:8px;padding:0.75rem 1rem;margin:0.75rem 0;color:#475569;font-size:0.825rem'>
        💡 <strong style='color:#64748b'>Pro tip:</strong> The more specific your question, the better DashAI's analysis.</div>""", unsafe_allow_html=True)


def _load_data():
    st.markdown('<div class="step-badge">Step 2 — Upload Data</div>', unsafe_allow_html=True)
    st.header("Upload Your Data")
    st.caption("CSV or Excel (multiple sheets auto-merged) • Up to 200MB")

    uploaded_files = st.file_uploader("Drop your files here", type=['csv','xlsx','xls'], accept_multiple_files=True, label_visibility="collapsed")

    if uploaded_files:
        dfs = []
        progress_bar = st.progress(0)
        for idx, f in enumerate(uploaded_files):
            try:
                if f.name.endswith('.csv'):
                    d = pd.read_csv(f); dfs.append(d)
                    st.caption(f"✓ {f.name}: {len(d):,} rows, {len(d.columns)} cols")
                else:
                    xls = pd.ExcelFile(f)
                    for sheet in xls.sheet_names:
                        d = pd.read_excel(xls, sheet_name=sheet)
                        if len(d) > 0:
                            dfs.append(d)
                            st.caption(f"✓ {f.name} → {sheet}: {len(d):,} rows, {len(d.columns)} cols")
            except Exception as e:
                st.warning(f"⚠️ Could not load {f.name}: {e}")
            progress_bar.progress((idx + 1) / len(uploaded_files))
        progress_bar.empty()
        if dfs:
            raw_df = dfs[0]
            if len(dfs) > 1:
                for other in dfs[1:]:
                    shared = list(set(raw_df.columns) & set(other.columns))
                    raw_df = raw_df.merge(other, on=shared, how='outer') if shared else pd.concat([raw_df, other], axis=1)
            prev = st.session_state.get('uploaded_df')
            if prev is None or list(raw_df.columns) != list(prev.columns):
                st.session_state['cleaned_df'] = None
                st.session_state['analyst_report'] = None
                st.session_state['quality_before'] = None
                st.session_state['quality_after'] = None
            st.session_state['uploaded_df'] = raw_df

    if st.session_state.get('cleaned_df') is not None:
        df = st.session_state['cleaned_df']
        st.success(f"✅ Using AI-cleaned data: **{len(df):,} rows × {len(df.columns)} columns**")
        return df, True
    elif st.session_state.get('uploaded_df') is not None:
        df = st.session_state['uploaded_df']
        if uploaded_files:
            st.success(f"✅ Loaded: **{len(df):,} rows × {len(df.columns)} columns**")
        return df, True
    else:
        return None, False


def main():
    st.markdown("""
    <div style='padding:1.25rem 0 1rem 0;border-bottom:1px solid #1e2035;margin-bottom:1.5rem'>
        <div style='display:flex;align-items:center;gap:0.6rem;margin-bottom:0.25rem'>
            <span style='font-size:1.4rem'>📊</span>
            <span style='font-size:1.4rem;font-weight:800;color:#f1f5f9;letter-spacing:-0.5px'>DashAI</span>
            <span style='background:#6366f115;border:1px solid #6366f130;border-radius:4px;padding:0.1rem 0.45rem;font-size:0.65rem;font-weight:700;color:#818cf8;text-transform:uppercase;letter-spacing:0.1em'>v11</span>
        </div>
        <p style='color:#64748b;font-size:0.8rem;margin:0;font-weight:500'>
        AI Dashboard Builder &nbsp;·&nbsp; Natural Language Queries &nbsp;·&nbsp; Auto Data Cleaning &nbsp;·&nbsp; PDF & PowerPoint Export
        </p>
    </div>
    """, unsafe_allow_html=True)

    client = get_claude_client()
    if not client:
        st.stop()

    if st.session_state['saved_dashboards']:
        with st.sidebar:
            st.markdown("### Saved Dashboards")
            for idx, saved in enumerate(st.session_state['saved_dashboards']):
                c1, c2 = st.columns([3, 1])
                with c1:
                    if st.button(f"📊 {saved['name']}", key=f"load_{idx}"):
                        st.session_state['dashboard_result'] = {'type': 'custom', 'data': saved['config']}
                        st.rerun()
                with c2:
                    if st.button("✕", key=f"del_{idx}"):
                        st.session_state['saved_dashboards'].pop(idx)
                        st.rerun()
                st.caption(saved['timestamp'])
            if st.button("Clear All"):
                st.session_state['saved_dashboards'] = []
                st.rerun()
            st.divider()

    has_bq = bool(st.session_state.get('business_question'))
    has_data = st.session_state.get('uploaded_df') is not None or st.session_state.get('cleaned_df') is not None
    has_analysis = st.session_state.get('analyst_report') is not None
    has_cleaned = st.session_state.get('cleaned_df') is not None

    def status_card(label, active, icon_on, icon_off):
        color = "#10b981" if active else "#1e2035"
        text_color = "#10b981" if active else "#475569"
        icon = icon_on if active else icon_off
        return f"<div style='background:#12131f;border:1px solid {color};border-radius:8px;padding:0.4rem 0.75rem;text-align:center;font-size:0.72rem;font-weight:600;color:{text_color}'>{icon} {label}</div>"

    sc1, sc2, sc3, sc4 = st.columns(4)
    with sc1: st.markdown(status_card("Question Set", has_bq, "✓", "○"), unsafe_allow_html=True)
    with sc2: st.markdown(status_card("Data Loaded", has_data, "✓", "○"), unsafe_allow_html=True)
    with sc3: st.markdown(status_card("Analysed", has_analysis, "✓", "○"), unsafe_allow_html=True)
    with sc4: st.markdown(status_card("Data Cleaned", has_cleaned, "✓", "○"), unsafe_allow_html=True)
    st.markdown("<br>", unsafe_allow_html=True)

    tab1, tab2, tab3, tab4, tab5 = st.tabs(["📋 1 · Setup", "🤖 2 · Analyse", "🔍 3 · Explore", "📊 4 · Dashboard", "📡 5 · Monitor"])

    with tab1:
        _render_step1_bq()
        st.divider()
        df, loaded = _load_data()
        if loaded and df is not None:
            with st.expander("Data Quality Check"):
                issues = validate_data(df)
                if not issues:
                    st.success("✅ No major issues detected!")
                else:
                    for issue in issues:
                        fn = st.warning if issue['severity'] == 'warning' else st.info
                        fn(f"**{issue['type']}**: {issue['message']}")
            with st.expander("Data Preview"):
                st.dataframe(df.head(10), use_container_width=True)

    with tab2:
        df = st.session_state.get('cleaned_df') if st.session_state.get('cleaned_df') is not None else st.session_state.get('uploaded_df')
        if df is None:
            st.info("👈 Go to **Setup** tab first to upload your data.")
        else:
            st.markdown('<div class="step-badge">Step 3 — AI Data Analyst</div>', unsafe_allow_html=True)
            st.header("AI Data Analyst")
            col1, col2 = st.columns([3, 1])
            with col1:
                st.markdown("<span style='color:#94a3b8;font-size:0.875rem'>📊 Data Profiling &nbsp;·&nbsp; 🔍 Issue Detection &nbsp;·&nbsp; 🧹 Auto-Cleaning &nbsp;·&nbsp; ⚡ Feature Engineering &nbsp;·&nbsp; 💡 Business Insights</span>", unsafe_allow_html=True)
            with col2:
                run_analyst = st.button("🤖 Run AI Analyst", type="primary", use_container_width=True, key="run_analyst_tab2")
                if st.session_state['analyst_report']:
                    if st.button("Re-run", use_container_width=True, key="rerun_analyst_tab2"):
                        st.session_state['analyst_report'] = None
                        st.session_state['cleaned_df'] = None
                        st.session_state['quality_before'] = None
                        st.session_state['quality_after'] = None
                        st.rerun()

            if run_analyst:
                with st.spinner("AI Analyst profiling, cleaning, and analyzing..."):
                    raw_df_for_score = st.session_state.get('uploaded_df', df)
                    report = run_ai_analyst(raw_df_for_score, client, business_question=st.session_state.get('business_question'))
                    if report:
                        st.session_state['analyst_report'] = report
                        st.session_state['quality_before'] = compute_quality_score(raw_df_for_score)
                        st.session_state['quality_after'] = None
                        st.rerun()

            if st.session_state.get('quality_before') is not None and st.session_state.get('quality_after') is not None:
                before = st.session_state['quality_before']
                after  = st.session_state['quality_after']
                delta  = after - before
                dc = '#10b981' if delta >= 0 else '#ef4444'
                di = '↑' if delta >= 0 else '↓'
                def _qcolor(v): return '#10b981' if v >= 80 else '#f59e0b' if v >= 60 else '#ef4444'
                st.markdown(f"""
<div style='background:#0d2318;border:1px solid #166534;border-radius:14px;padding:1rem 1.5rem;margin:0.5rem 0 1rem 0'>
  <div style='font-size:0.68rem;font-weight:700;color:#10b981;text-transform:uppercase;letter-spacing:0.1em;margin-bottom:0.75rem'>✅ Data Quality — Before vs After Cleaning</div>
  <div style='display:flex;align-items:center;gap:2.5rem;flex-wrap:wrap'>
    <div style='text-align:center;min-width:90px'>
      <div style='font-size:0.65rem;color:#64748b;font-weight:600;text-transform:uppercase;letter-spacing:0.06em;margin-bottom:0.25rem'>Before</div>
      <div style='font-size:2.4rem;font-weight:800;color:{_qcolor(before)}'>{before}%</div>
    </div>
    <div style='text-align:center;flex:1;min-width:80px'>
      <div style='font-size:2.5rem;color:#2d3050'>→</div>
      <div style='font-size:1.05rem;font-weight:700;color:{dc}'>{di} {abs(delta)} pts</div>
    </div>
    <div style='text-align:center;min-width:90px'>
      <div style='font-size:0.65rem;color:#64748b;font-weight:600;text-transform:uppercase;letter-spacing:0.06em;margin-bottom:0.25rem'>After</div>
      <div style='font-size:2.4rem;font-weight:800;color:{_qcolor(after)}'>{after}%</div>
    </div>
    <div style='flex:2;min-width:160px;padding-left:1rem;border-left:1px solid #166534'>
      <div style='font-size:0.75rem;color:#86efac;font-weight:600;margin-bottom:0.25rem'>What improved:</div>
      <div style='font-size:0.78rem;color:#64748b'>• Missing values imputed</div>
      <div style='font-size:0.78rem;color:#64748b'>• Duplicates removed</div>
      <div style='font-size:0.78rem;color:#64748b'>• Types fixed, features engineered</div>
    </div>
  </div>
</div>""", unsafe_allow_html=True)

            if st.session_state['analyst_report']:
                render_analyst_report(st.session_state['analyst_report'], df, client)
                if st.session_state['cleaned_df'] is not None:
                    df = st.session_state['cleaned_df']

            st.divider()
            st.markdown('<div class="step-badge">Step 4 — Recommendations</div>', unsafe_allow_html=True)
            st.header("AI Recommendations")
            if st.button("🎯 Generate Recommendations", type="primary", use_container_width=True, key="gen_recs_tab2"):
                with st.spinner("Synthesizing recommendations..."):
                    rec = run_recommendations_engine(df, client, analyst_report=st.session_state.get('analyst_report'), business_question=st.session_state.get('business_question'))
                    if rec:
                        st.session_state['recommendations_result'] = rec
                        st.rerun()

            if st.session_state.get('recommendations_result'):
                render_recommendations(st.session_state['recommendations_result'], df, client)

    with tab3:
        df = st.session_state.get('cleaned_df') if st.session_state.get('cleaned_df') is not None else st.session_state.get('uploaded_df')
        if df is None:
            st.info("👈 Go to **Setup** tab first to upload your data.")
        else:
            st.markdown('<div class="step-badge">Step 5 — Natural Language Query</div>', unsafe_allow_html=True)
            st.header("Ask Questions")
            st.caption("Ask anything about your data in plain English — get real answers + charts")
            c1, c2 = st.columns([4, 1])
            with c1:
                question = st.text_input("Your question", placeholder="e.g., Which state has the highest churn rate?", label_visibility="collapsed", key="nl_question_tab3")
            with c2:
                ask_btn = st.button("Ask →", type="primary", use_container_width=True, key="ask_btn_tab3")

            if ask_btn and question:
                with st.spinner("Analyzing..."):
                    result = natural_language_query(df, question, client)
                if result:
                    st.markdown(f"<div style='color:#e2e8f0;font-size:0.9rem;margin:0.5rem 0'>{result.get('answer','')}</div>", unsafe_allow_html=True)
                    if result.get('insight'):
                        st.markdown(f"""<div style='background:#12131f;border-left:2px solid #6366f1;border-radius:0 8px 8px 0;padding:0.5rem 0.9rem;margin:0.4rem 0;color:#6366f1;font-size:0.825rem'>💡 {result['insight']}</div>""", unsafe_allow_html=True)
                    code = result.get('code', '')
                    if code:
                        out = safe_exec(code, df)
                        if out['result'] is not None and hasattr(out['result'], 'update_layout'):
                            fig = style_fig(out['result'])
                            st.plotly_chart(fig, use_container_width=True)
                        elif out['error']:
                            st.warning(f"Chart error: {out['error'][:100]}")
                        with st.expander("View Code"):
                            st.code(code, language='python')
                    st.session_state['query_history'].append({'q': question, 'time': datetime.now().strftime('%H:%M')})

            if st.session_state['query_history']:
                with st.expander("Query History"):
                    for q in reversed(st.session_state['query_history'][-5:]):
                        st.markdown(f"<span style='color:#64748b;font-size:0.75rem'>{q['time']}</span> <span style='color:#94a3b8;font-size:0.85rem'>{q['q']}</span>", unsafe_allow_html=True)

    with tab4:
        df = st.session_state.get('cleaned_df') if st.session_state.get('cleaned_df') is not None else st.session_state.get('uploaded_df')
        if df is None:
            st.info("👈 Go to **Setup** tab first to upload your data.")
        else:
            st.markdown('<div class="step-badge">Step 6 — Custom Dashboard</div>', unsafe_allow_html=True)
            st.header("Generate Custom Dashboard")
            st.caption("Use a template or describe exactly what you want")

            t_cols = st.columns(6)
            for i, (name, tpl) in enumerate(PROMPT_TEMPLATES.items()):
                with t_cols[i]:
                    if st.button(f"{tpl['emoji']} {name}", use_container_width=True, key=f"tpl_btn_{name}"):
                        st.session_state['prompt_template'] = tpl['prompt']
                        st.session_state['_dash_desc_textarea'] = tpl['prompt']
                        st.rerun()

            if '_dash_desc_textarea' not in st.session_state:
                st.session_state['_dash_desc_textarea'] = st.session_state.get('prompt_template', '')

            custom_prompt = st.text_area(
                "Dashboard Description",
                placeholder="e.g., 'Executive overview with revenue KPIs, plan mix donut chart, and state map'",
                height=100,
                key="_dash_desc_textarea",
            )
            st.session_state['prompt_template'] = custom_prompt

            c1, c2 = st.columns([3, 1])
            with c1:
                if st.button("🚀 Generate Dashboard", type="primary", use_container_width=True, key="btn_generate_dashboard"):
                    if custom_prompt and custom_prompt.strip():
                        with st.spinner("Generating dashboard..."):
                            prog = st.progress(0)
                            for i in range(100):
                                time.sleep(0.008)
                                prog.progress(i + 1)
                            result = generate_dashboard(df, custom_prompt, client)
                            prog.empty()
                            if result:
                                st.session_state['dashboard_result'] = {'type': 'custom', 'data': result}
                                st.rerun()
                    else:
                        st.warning("Enter a description or pick a template first!")
            with c2:
                if st.button("Clear", use_container_width=True, key="btn_clear_dashboard"):
                    st.session_state['prompt_template'] = ''
                    st.session_state['_dash_desc_textarea'] = ''
                    st.rerun()

            st.divider()
            st.markdown('<div class="step-badge">Step 7 — Preset Role Dashboards</div>', unsafe_allow_html=True)
            st.header("Or Choose a Preset Dashboard")
            if st.button("📊 Generate 5 Role-Based Dashboards", use_container_width=True, key="btn_gen_presets"):
                with st.spinner("Generating 5 role-based dashboards..."):
                    prog = st.progress(0)
                    for i in range(100):
                        time.sleep(0.015)
                        prog.progress(i + 1)
                    result = generate_preset_dashboards(df, client)
                    prog.empty()
                    if result:
                        st.session_state['dashboard_result'] = {'type': 'presets', 'data': result}
                        st.rerun()

            dash_result = st.session_state.get('dashboard_result')
            if dash_result:
                if dash_result['type'] == 'custom':
                    st.divider()
                    render_dashboard(df, dash_result['data'])
                    if st.button("← Generate New Dashboard", use_container_width=True, key="btn_new_dashboard"):
                        st.session_state['dashboard_result'] = None
                        st.session_state['prompt_template'] = ''
                        st.rerun()

                elif dash_result['type'] == 'presets':
                    dashboards = dash_result['data'].get('dashboards', [])
                    st.divider()
                    st.success(f"✅ {len(dashboards)} dashboards ready — click one to open it")
                    cols = st.columns(min(len(dashboards), 5))
                    for i, d in enumerate(dashboards):
                        with cols[i % 5]:
                            st.markdown(f"""<div style='background:#12131f;border:1px solid #1e2035;border-radius:12px;padding:1rem;text-align:center;margin-bottom:0.5rem'>
                            <div style='font-size:1.5rem;margin-bottom:0.3rem'>{d.get("icon","📊")}</div>
                            <div style='color:#e2e8f0;font-weight:700;font-size:0.825rem;margin-bottom:0.15rem'>{d.get("name","")}</div>
                            <div style='color:#6366f1;font-size:0.7rem;font-weight:600;margin-bottom:0.25rem'>{d.get("audience","")}</div>
                            <div style='color:#64748b;font-size:0.72rem;line-height:1.3'>{d.get("description","")}</div></div>""", unsafe_allow_html=True)
                            if st.button("Open →", key=f"preset_{i}", use_container_width=True):
                                st.session_state['dashboard_result'] = {'type': 'custom', 'data': d}
                                st.rerun()
                    if st.button("Regenerate All", use_container_width=True, key="btn_regen_presets"):
                        st.session_state['dashboard_result'] = None
                        st.rerun()

    with tab5:
        df_mon = st.session_state.get('cleaned_df') if st.session_state.get('cleaned_df') is not None else st.session_state.get('uploaded_df')
        if df_mon is None:
            st.info("👈 Go to **Setup** tab first to upload your data.")
        else:
            render_monitor_tab(df_mon, client)


if __name__ == "__main__":
    main()
